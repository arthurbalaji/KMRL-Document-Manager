import os
import logging
import json
import traceback
import hashlib
from datetime import datetime

# Valid KMRL roles - must match backend Role enum exactly
VALID_KMRL_ROLES = ['LEADERSHIP', 'HR', 'FINANCE', 'ENGINEER', 'ADMIN']

def validate_and_filter_roles(roles_list):
    """Filter AI-suggested roles to only include valid KMRL roles."""
    if not isinstance(roles_list, list):
        return ['LEADERSHIP']
    
    original_roles = roles_list.copy()
    valid_roles = [role for role in roles_list if role in VALID_KMRL_ROLES]
    
    # Log any invalid roles that were filtered out
    invalid_roles = [role for role in original_roles if role not in VALID_KMRL_ROLES]
    if invalid_roles:
        logger.warning(f"Filtered out invalid roles from AI response: {invalid_roles}")
        logger.info(f"Valid roles after filtering: {valid_roles}")
    
    return valid_roles if valid_roles else ['LEADERSHIP']

def clamp01(x: float) -> float:
    try:
        # Ensure confidence is never zero - minimum of 0.1 (10%)
        return max(0.1, min(1.0, float(x)))
    except Exception:
        return 0.75

def normalize_analysis(analysis: dict) -> dict:
    """Ensure analysis has recommended_roles with a non-null confidence in [0.1,1] and reasonable defaults."""
    logger.info(f"normalize_analysis called with: {type(analysis)}")
    try:
        if not isinstance(analysis, dict):
            logger.warning("Analysis is not a dict, using defaults")
            return {
                'summary_en': '',
                'summary_ml': '',
                'tags': [],
                'primary_language': 'en',
                'sensitivity_level': 'MEDIUM',
                'document_type': 'other',
                'recommended_roles': {
                    'roles': ['LEADERSHIP'],
                    'confidence': 0.75,
                    'reasoning': 'Defaulted due to invalid analysis structure'
                },
                'retention_recommendation': {'days': 1825, 'reason': 'Default policy'},
                'analysis_metadata': {}
            }

        rec = analysis.get('recommended_roles') or {}
        roles = rec.get('roles') or ['LEADERSHIP']
        # Validate and filter roles to only include valid KMRL roles
        roles = validate_and_filter_roles(roles)
        
        conf = rec.get('confidence', 0.75)
        logger.info(f"Original confidence from AI: {conf} (type: {type(conf)})")
        
        try:
            conf = clamp01(conf)
            logger.info(f"After clamp01: {conf}")
        except Exception as e:
            logger.warning(f"Error in clamp01: {e}")
            conf = 0.75

        # If model returned 0 exactly, lift slightly to avoid zero downstream visuals
        if conf == 0.0:
            logger.info("Confidence was 0.0, setting to 0.15")
            conf = 0.15

        reasoning = rec.get('reasoning') or 'Combined AI and heuristic analysis.'
        analysis['recommended_roles'] = {
            'roles': roles,
            'confidence': conf,
            'reasoning': reasoning
        }

        # Default sensitivity/document_type if missing
        if not analysis.get('sensitivity_level'):
            analysis['sensitivity_level'] = 'MEDIUM'
        if not analysis.get('document_type'):
            analysis['document_type'] = 'other'

        return analysis
    except Exception as e:
        logger.warning(f"normalize_analysis failed: {e}")
        return {
            'summary_en': analysis.get('summary_en', '') if isinstance(analysis, dict) else '',
            'summary_ml': analysis.get('summary_ml', '') if isinstance(analysis, dict) else '',
            'tags': analysis.get('tags', []) if isinstance(analysis, dict) else [],
            'primary_language': analysis.get('primary_language', 'en') if isinstance(analysis, dict) else 'en',
            'sensitivity_level': analysis.get('sensitivity_level', 'MEDIUM') if isinstance(analysis, dict) else 'MEDIUM',
            'document_type': analysis.get('document_type', 'other') if isinstance(analysis, dict) else 'other',
            'recommended_roles': {
                'roles': ['LEADERSHIP'],
                'confidence': 0.75,
                'reasoning': 'Defaulted due to exception in normalization'
            },
            'retention_recommendation': {'days': 1825, 'reason': 'Default policy'},
            'analysis_metadata': {}
        }
from typing import Dict, List, Tuple, Optional
import base64
import io
import tempfile
import uuid
import csv
import chardet

import requests
from flask import Flask, request, jsonify, send_file, make_response
from flask_cors import CORS
import psycopg2
from psycopg2.extras import RealDictCursor
from dotenv import load_dotenv
import pytesseract
from PIL import Image
import PyPDF2
from docx import Document
import numpy as np
import openpyxl
import xlrd
import pandas as pd
from pptx import Presentation

# Fix for PIL.Image.ANTIALIAS deprecation in newer versions
try:
    # Try to use the old ANTIALIAS constant
    Image.ANTIALIAS
except AttributeError:
    # If it doesn't exist, use LANCZOS as replacement
    Image.ANTIALIAS = Image.Resampling.LANCZOS
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from langdetect import detect, DetectorFactory
import re
import cv2
import easyocr
from googletrans import Translator
import pdf2image
import pdfplumber
import fitz  # PyMuPDF
from collections import defaultdict

# Ensure consistent language detection
DetectorFactory.seed = 0

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize Flask app
app = Flask(__name__)
CORS(app)

# Configure DeepSeek V2
deepseek_api_key = os.getenv('DEEPSEEK_API_KEY')
if not deepseek_api_key:
    logger.error("DEEPSEEK_API_KEY not found in environment variables")
    raise ValueError("DEEPSEEK_API_KEY is required")

logger.info(f"Configuring DeepSeek V2 with API key: {deepseek_api_key[:10]}...")

# Configure Tesseract path
TESSERACT_CMD = os.getenv('TESSERACT_CMD', '/usr/bin/tesseract')

# DeepSeek API configuration
DEEPSEEK_BASE_URL = "https://api.deepseek.com/v1/chat/completions"
DEEPSEEK_HEADERS = {
    "Authorization": f"Bearer {deepseek_api_key}",
    "Content-Type": "application/json"
}

def call_deepseek_api(messages, max_tokens=4000, temperature=0.3):
    """Call DeepSeek API using direct HTTP requests"""
    try:
        payload = {
            "model": "deepseek-chat",
            "messages": messages,
            "max_tokens": max_tokens,
            "temperature": temperature
        }
        
        response = requests.post(
            DEEPSEEK_BASE_URL,
            headers=DEEPSEEK_HEADERS,
            json=payload,
            timeout=60
        )
        
        if response.status_code == 200:
            data = response.json()
            return data["choices"][0]["message"]["content"]
        else:
            logger.error(f"DeepSeek API error: {response.status_code} - {response.text}")
            return None
            
    except Exception as e:
        logger.error(f"DeepSeek API call failed: {str(e)}")
        return None

def call_deepseek_api_with_caching(cache_key, messages, max_tokens=4000, temperature=0.3, cache_duration=3600):
    """Enhanced DeepSeek API call with caching to reduce redundant API usage"""
    import time
    import hashlib
    
    # Simple in-memory cache (in production, use Redis or similar)
    if not hasattr(call_deepseek_api_with_caching, 'cache'):
        call_deepseek_api_with_caching.cache = {}
    
    # Create cache key from messages and parameters
    if cache_key is None:
        cache_content = json.dumps(messages) + str(max_tokens) + str(temperature)
        cache_key = hashlib.md5(cache_content.encode()).hexdigest()
    
    # Check cache
    cache = call_deepseek_api_with_caching.cache
    current_time = time.time()
    
    if cache_key in cache:
        cached_item = cache[cache_key]
        if current_time - cached_item['timestamp'] < cache_duration:
            logger.info(f"Cache hit for key: {cache_key[:10]}...")
            return cached_item['result']
        else:
            # Remove expired cache
            del cache[cache_key]
    
    # Make API call
    result = call_deepseek_api(messages, max_tokens, temperature)
    
    # Cache the result if successful
    if result:
        cache[cache_key] = {
            'result': result,
            'timestamp': current_time
        }
        
        # Simple cache cleanup - remove old entries if cache gets too large
        if len(cache) > 100:
            oldest_key = min(cache.keys(), key=lambda k: cache[k]['timestamp'])
            del cache[oldest_key]
    
    return result

def deepseek_enhanced_language_detection(text: str) -> Dict:
    """Use DeepSeek API for enhanced language detection with confidence scores"""
    try:
        if len(text.strip()) < 10:
            return {'language': 'en', 'confidence': 0.5, 'languages': ['en']}
        
        prompt = f"""
        Analyze the following text and detect the languages present:
        
        Text: {text[:1000]}  # Limit text for efficiency
        
        Return a JSON response with this exact structure:
        {{
            "primary_language": "en|ml|hi|ta|te|...",
            "confidence": 0.95,
            "all_languages": ["en", "ml"],
            "language_percentages": {{"en": 60, "ml": 40}},
            "is_multilingual": true
        }}
        
        Guidelines:
        - primary_language: ISO 639-1 code for the dominant language
        - confidence: 0.0-1.0 confidence in detection
        - all_languages: list of all detected languages
        - language_percentages: approximate percentage of each language
        - is_multilingual: true if more than one language detected
        """
        
        response_text = optimized_deepseek_call(
            "language_detection",
            text[:500],
            prompt,
            {'max_tokens': 500, 'temperature': 0.1}
        )
        
        if response_text:
            # Clean and parse JSON
            response_text = response_text.strip()
            if response_text.startswith('```json'):
                response_text = response_text[7:]
            if response_text.endswith('```'):
                response_text = response_text[:-3]
            
            result = json.loads(response_text)
            return {
                'language': result.get('primary_language', 'en'),
                'confidence': result.get('confidence', 0.8),
                'languages': result.get('all_languages', ['en']),
                'percentages': result.get('language_percentages', {}),
                'is_multilingual': result.get('is_multilingual', False)
            }
        else:
            raise Exception("DeepSeek API returned no response")
            
    except Exception as e:
        logger.warning(f"Enhanced language detection failed: {str(e)}, falling back to local detection")
        # Fallback to local detection
        try:
            detected = detect(text)
            return {
                'language': detected,
                'confidence': 0.7,
                'languages': [detected],
                'percentages': {detected: 100},
                'is_multilingual': False
            }
        except:
            return {
                'language': 'en',
                'confidence': 0.5,
                'languages': ['en'],
                'percentages': {'en': 100},
                'is_multilingual': False
            }

def deepseek_enhanced_translation(text: str, target_lang: str = 'en', source_lang: str = None) -> Dict:
    """Use DeepSeek API for context-aware translation"""
    try:
        if not text.strip():
            return {'translated_text': text, 'confidence': 1.0, 'source_language': 'unknown'}
        
        # Auto-detect source language if not provided
        if not source_lang:
            lang_info = deepseek_enhanced_language_detection(text)
            source_lang = lang_info['language']
        
        # Don't translate if already in target language
        if source_lang == target_lang:
            return {'translated_text': text, 'confidence': 1.0, 'source_language': source_lang}
        
        prompt = f"""
        Translate the following text from {source_lang} to {target_lang}.
        Preserve the meaning, context, and any technical terms accurately.
        This is for a metro rail document management system.
        
        Source text:
        {text[:2000]}
        
        Return only the translated text without any explanations or formatting.
        """
        
        cache_key = f"translate_{source_lang}_{target_lang}_{hashlib.md5(text[:500].encode()).hexdigest()}"
        translated = call_deepseek_api_with_caching(
            cache_key,
            [{"role": "user", "content": prompt}],
            max_tokens=len(text.split()) * 2,  # Estimate tokens needed
            temperature=0.1
        )
        
        if translated:
            return {
                'translated_text': translated.strip(),
                'confidence': 0.9,
                'source_language': source_lang
            }
        else:
            raise Exception("DeepSeek translation failed")
            
    except Exception as e:
        logger.warning(f"DeepSeek translation failed: {str(e)}, falling back to Google Translate")
        # Fallback to existing Google Translate
        try:
            detected = translator.detect(text)
            if detected.lang == target_lang:
                return {'translated_text': text, 'confidence': 1.0, 'source_language': detected.lang}
            
            translated = translator.translate(text, dest=target_lang)
            return {
                'translated_text': translated.text,
                'confidence': 0.7,
                'source_language': detected.lang
            }
        except:
            return {'translated_text': text, 'confidence': 0.5, 'source_language': 'unknown'}

def deepseek_enhanced_embeddings(text: str) -> List[float]:
    """Generate semantic embeddings using DeepSeek API analysis"""
    try:
        if not text.strip():
            return [0.0] * 512
        
        prompt = f"""
        Analyze the following text and extract semantic features for document similarity matching.
        Create a comprehensive semantic profile focusing on:
        - Key concepts and topics
        - Document type and purpose
        - Domain-specific terminology
        - Sentiment and tone
        - Structural elements
        
        Text: {text[:3000]}
        
        Return a JSON with semantic features as a list of 50 normalized float values (0.0-1.0):
        {{
            "semantic_vector": [0.23, 0.45, 0.67, ...]
        }}
        
        Each value should represent semantic strength in different dimensions.
        """
        
        cache_key = f"embeddings_{hashlib.md5(text[:1000].encode()).hexdigest()}"
        response_text = call_deepseek_api_with_caching(
            cache_key,
            [{"role": "user", "content": prompt}],
            max_tokens=1000,
            temperature=0.1
        )
        
        if response_text:
            # Parse semantic vector
            response_text = response_text.strip()
            if response_text.startswith('```json'):
                response_text = response_text[7:]
            if response_text.endswith('```'):
                response_text = response_text[:-3]
            
            result = json.loads(response_text)
            semantic_vector = result.get('semantic_vector', [])
            
            # Ensure we have 512 dimensions by extending/truncating
            if len(semantic_vector) < 512:
                # Pad with computed features from local TF-IDF
                local_embeddings = generate_local_embeddings(text)
                semantic_vector.extend(local_embeddings[len(semantic_vector):])
            
            return semantic_vector[:512]
        else:
            raise Exception("DeepSeek embeddings failed")
            
    except Exception as e:
        logger.warning(f"DeepSeek embeddings failed: {str(e)}, falling back to local TF-IDF")
        return generate_local_embeddings(text)

def generate_local_embeddings(text: str) -> List[float]:
    """Fallback local embedding generation using TF-IDF"""
    try:
        vectorizer = TfidfVectorizer(max_features=512, stop_words='english')
        tfidf_matrix = vectorizer.fit_transform([text])
        embeddings = tfidf_matrix.toarray()[0].tolist()
        return embeddings
    except Exception as e:
        logger.error(f"Local embedding generation error: {str(e)}")
        return [0.0] * 512

def deepseek_batch_processing(batch_requests: List[Dict]) -> List[Dict]:
    """Process multiple DeepSeek API requests in optimized batches"""
    try:
        results = []
        
        # Group similar requests for batch processing
        grouped_requests = {
            'language_detection': [],
            'translation': [],
            'ocr_enhancement': [],
            'embeddings': [],
            'analysis': []
        }
        
        # Group requests by type
        for i, request in enumerate(batch_requests):
            request['original_index'] = i
            request_type = request.get('type')
            if request_type in grouped_requests:
                grouped_requests[request_type].append(request)
            else:
                # Process unknown types individually
                grouped_requests['analysis'].append(request)
        
        # Process each group
        for request_type, requests in grouped_requests.items():
            if not requests:
                continue
                
            if request_type == 'language_detection':
                # Batch language detection
                batch_results = batch_language_detection([r['text'] for r in requests])
                for i, result in enumerate(batch_results):
                    results.append({
                        'index': requests[i]['original_index'],
                        'result': result
                    })
                    
            elif request_type == 'translation':
                # Batch translation
                batch_results = batch_translation(requests)
                for i, result in enumerate(batch_results):
                    results.append({
                        'index': requests[i]['original_index'],
                        'result': result
                    })
                    
            else:
                # Process individually for other types
                for request in requests:
                    if request_type == 'ocr_enhancement':
                        result = deepseek_ocr_enhancement(request['text'], request.get('context', ''))
                    elif request_type == 'embeddings':
                        result = deepseek_enhanced_embeddings(request['text'])
                    else:
                        result = {'error': 'Unknown request type'}
                    
                    results.append({
                        'index': request['original_index'],
                        'result': result
                    })
        
        # Sort results by original index
        results.sort(key=lambda x: x['index'])
        return [r['result'] for r in results]
        
    except Exception as e:
        logger.error(f"Batch processing error: {str(e)}")
        return [{'error': str(e)} for _ in batch_requests]

def batch_language_detection(texts: List[str]) -> List[Dict]:
    """Batch process language detection for multiple texts"""
    try:
        if not texts:
            return []
        
        # Combine texts for single API call
        combined_prompt = f"""
        Detect languages for the following {len(texts)} text samples.
        Return a JSON array with detection results for each text:
        
        Texts:
        """
        
        for i, text in enumerate(texts):
            combined_prompt += f"\n{i+1}. {text[:300]}...\n"
        
        combined_prompt += """
        
        Return JSON array format:
        [
            {
                "text_index": 1,
                "primary_language": "en",
                "confidence": 0.95,
                "all_languages": ["en"],
                "language_percentages": {"en": 100},
                "is_multilingual": false
            },
            ...
        ]
        """
        
        response_text = call_deepseek_api_with_caching(
            None,  # No specific cache key for batch
            [{"role": "user", "content": combined_prompt}],
            max_tokens=len(texts) * 200,
            temperature=0.1
        )
        
        if response_text:
            # Parse batch results
            response_text = response_text.strip()
            if response_text.startswith('```json'):
                response_text = response_text[7:]
            if response_text.endswith('```'):
                response_text = response_text[:-3]
            
            batch_results = json.loads(response_text)
            
            # Convert to expected format
            results = []
            for i, text in enumerate(texts):
                # Find corresponding result
                result_item = None
                for item in batch_results:
                    if item.get('text_index') == i + 1:
                        result_item = item
                        break
                
                if result_item:
                    results.append({
                        'language': result_item.get('primary_language', 'en'),
                        'confidence': result_item.get('confidence', 0.8),
                        'languages': result_item.get('all_languages', ['en']),
                        'percentages': result_item.get('language_percentages', {}),
                        'is_multilingual': result_item.get('is_multilingual', False)
                    })
                else:
                    # Fallback for missing result
                    results.append(deepseek_enhanced_language_detection(text))
            
            return results
        else:
            raise Exception("Batch language detection API call failed")
            
    except Exception as e:
        logger.warning(f"Batch language detection failed: {str(e)}")
        # Fallback to individual processing
        return [deepseek_enhanced_language_detection(text) for text in texts]

def batch_translation(translation_requests: List[Dict]) -> List[Dict]:
    """Batch process translation requests"""
    try:
        if not translation_requests:
            return []
        
        # Group by target language for efficiency
        lang_groups = {}
        for i, req in enumerate(translation_requests):
            target_lang = req.get('target_lang', 'en')
            if target_lang not in lang_groups:
                lang_groups[target_lang] = []
            lang_groups[target_lang].append((i, req))
        
        results = [None] * len(translation_requests)
        
        # Process each language group
        for target_lang, requests in lang_groups.items():
            if len(requests) == 1:
                # Single request - process normally
                idx, req = requests[0]
                result = deepseek_enhanced_translation(req['text'], target_lang, req.get('source_lang'))
                results[idx] = result
            else:
                # Multiple requests - batch them
                combined_prompt = f"""
                Translate the following {len(requests)} texts to {target_lang}.
                Preserve meaning and context for metro rail documentation.
                
                Texts to translate:
                """
                
                for i, (idx, req) in enumerate(requests):
                    combined_prompt += f"\n{i+1}. {req['text'][:500]}...\n"
                
                combined_prompt += f"""
                
                Return JSON array with translations:
                [
                    {{"text_index": 1, "translated_text": "...", "confidence": 0.9}},
                    ...
                ]
                """
                
                response_text = call_deepseek_api_with_caching(
                    None,
                    [{"role": "user", "content": combined_prompt}],
                    max_tokens=sum(len(req['text'].split()) for _, req in requests) * 2,
                    temperature=0.1
                )
                
                if response_text:
                    response_text = response_text.strip()
                    if response_text.startswith('```json'):
                        response_text = response_text[7:]
                    if response_text.endswith('```'):
                        response_text = response_text[:-3]
                    
                    batch_results = json.loads(response_text)
                    
                    for i, (original_idx, req) in enumerate(requests):
                        result_item = None
                        for item in batch_results:
                            if item.get('text_index') == i + 1:
                                result_item = item
                                break
                        
                        if result_item:
                            results[original_idx] = {
                                'translated_text': result_item.get('translated_text', req['text']),
                                'confidence': result_item.get('confidence', 0.8),
                                'source_language': req.get('source_lang', 'auto')
                            }
                        else:
                            # Fallback
                            results[original_idx] = deepseek_enhanced_translation(
                                req['text'], target_lang, req.get('source_lang')
                            )
                else:
                    # Fallback for failed batch
                    for original_idx, req in requests:
                        results[original_idx] = deepseek_enhanced_translation(
                            req['text'], target_lang, req.get('source_lang')
                        )
        
        return results
        
    except Exception as e:
        logger.warning(f"Batch translation failed: {str(e)}")
        # Fallback to individual processing
        return [deepseek_enhanced_translation(req['text'], req.get('target_lang', 'en'), req.get('source_lang')) 
                for req in translation_requests]

def deepseek_ocr_enhancement(ocr_text: str, image_context: str = "") -> Dict:
    """Use DeepSeek API to enhance and correct OCR output"""
    try:
        if not ocr_text.strip():
            return {'enhanced_text': ocr_text, 'confidence': 0.0, 'corrections_made': []}
        
        prompt = f"""
        The following text was extracted using OCR and may contain errors. 
        Please correct any obvious OCR mistakes, improve formatting, and enhance readability.
        This is from a document in a metro rail system context.
        
        Image context: {image_context}
        
        OCR Text:
        {ocr_text}
        
        Return a JSON response:
        {{
            "corrected_text": "Enhanced and corrected text",
            "confidence": 0.95,
            "corrections": ["Original -> Corrected", ...],
            "detected_entities": ["Metro", "Station", ...]
        }}
        
        Guidelines:
        - Fix obvious character recognition errors (0/O, 1/l/I, etc.)
        - Improve spacing and formatting
        - Preserve technical terms and proper nouns
        - Don't change meaning, only improve clarity
        """
        
        cache_key = f"ocr_enhance_{hashlib.md5(ocr_text[:500].encode()).hexdigest()}"
        response_text = call_deepseek_api_with_caching(
            cache_key,
            [{"role": "user", "content": prompt}],
            max_tokens=len(ocr_text.split()) * 2,
            temperature=0.1
        )
        
        if response_text:
            # Clean and parse JSON with better error handling
            response_text = response_text.strip()
            
            # Remove markdown code blocks
            if response_text.startswith('```json'):
                response_text = response_text[7:]
            if response_text.startswith('```'):
                response_text = response_text[3:]
            if response_text.endswith('```'):
                response_text = response_text[:-3]
            
            # Try to extract JSON if it's embedded in other text
            json_start = response_text.find('{')
            json_end = response_text.rfind('}')
            
            if json_start != -1 and json_end != -1 and json_end > json_start:
                response_text = response_text[json_start:json_end+1]
            
            try:
                result = json.loads(response_text)
                return {
                    'enhanced_text': result.get('corrected_text', ocr_text),
                    'confidence': result.get('confidence', 0.8),
                    'corrections_made': result.get('corrections', []),
                    'entities': result.get('detected_entities', [])
                }
            except json.JSONDecodeError as je:
                logger.warning(f"JSON parsing failed for OCR enhancement: {str(je)}")
                logger.debug(f"Raw response (first 200 chars): {response_text[:200]}")
                
                # Try to extract just the corrected text if JSON parsing fails
                lines = response_text.split('\n')
                enhanced_text = ocr_text
                
                for line in lines:
                    if 'corrected_text' in line.lower() or 'enhanced' in line.lower():
                        # Try to extract text after colon
                        if ':' in line:
                            potential_text = line.split(':', 1)[1].strip().strip('"').strip("'")
                            if len(potential_text) > 10:  # Reasonable text length
                                enhanced_text = potential_text
                                break
                
                return {
                    'enhanced_text': enhanced_text,
                    'confidence': 0.6,
                    'corrections_made': [],
                    'entities': []
                }
        else:
            raise Exception("DeepSeek OCR enhancement returned empty response")
            
    except Exception as e:
        logger.warning(f"DeepSeek OCR enhancement failed: {str(e)}")
        return {'enhanced_text': ocr_text, 'confidence': 0.6, 'corrections_made': [], 'entities': []}

# API Usage Optimization - Rate limiting and intelligent caching
import time
from collections import defaultdict

class DeepSeekAPIOptimizer:
    def __init__(self):
        self.request_counts = defaultdict(int)
        self.request_timestamps = defaultdict(list)
        self.cache = {}
        self.max_requests_per_minute = 50  # Adjust based on API limits
        self.cache_ttl = 3600  # 1 hour cache TTL
        
    def can_make_request(self) -> bool:
        """Check if we can make a request without hitting rate limits"""
        current_time = time.time()
        minute_ago = current_time - 60
        
        # Clean old timestamps
        self.request_timestamps['minute'] = [
            ts for ts in self.request_timestamps['minute'] 
            if ts > minute_ago
        ]
        
        return len(self.request_timestamps['minute']) < self.max_requests_per_minute
    
    def record_request(self):
        """Record a new API request"""
        current_time = time.time()
        self.request_timestamps['minute'].append(current_time)
        self.request_counts['total'] += 1
    
    def get_cache_key(self, operation: str, content: str, params: dict = None) -> str:
        """Generate cache key for operation"""
        params_str = json.dumps(params or {}, sort_keys=True)
        content_hash = hashlib.md5(content[:1000].encode()).hexdigest()
        return f"{operation}_{content_hash}_{hashlib.md5(params_str.encode()).hexdigest()}"
    
    def get_cached_result(self, cache_key: str):
        """Get cached result if valid"""
        if cache_key in self.cache:
            cached_item = self.cache[cache_key]
            if time.time() - cached_item['timestamp'] < self.cache_ttl:
                return cached_item['result']
            else:
                del self.cache[cache_key]
        return None
    
    def cache_result(self, cache_key: str, result):
        """Cache a result"""
        self.cache[cache_key] = {
            'result': result,
            'timestamp': time.time()
        }
        
        # Cleanup old cache entries
        if len(self.cache) > 1000:
            oldest_key = min(self.cache.keys(), key=lambda k: self.cache[k]['timestamp'])
            del self.cache[oldest_key]
    
    def get_usage_stats(self) -> Dict:
        """Get API usage statistics"""
        current_time = time.time()
        minute_ago = current_time - 60
        hour_ago = current_time - 3600
        
        minute_requests = len([ts for ts in self.request_timestamps['minute'] if ts > minute_ago])
        total_requests = self.request_counts['total']
        cache_hits = len(self.cache)
        
        return {
            'requests_last_minute': minute_requests,
            'total_requests': total_requests,
            'cache_entries': cache_hits,
            'rate_limit_remaining': max(0, self.max_requests_per_minute - minute_requests)
        }

# Initialize API optimizer
api_optimizer = DeepSeekAPIOptimizer()

def optimized_deepseek_call(operation: str, content: str, prompt: str, params: dict = None) -> str:
    """Optimized DeepSeek API call with caching and rate limiting"""
    try:
        # Generate cache key
        cache_key = api_optimizer.get_cache_key(operation, content, params)
        
        # Try cache first
        cached_result = api_optimizer.get_cached_result(cache_key)
        if cached_result:
            logger.info(f"Cache hit for {operation}")
            return cached_result
        
        # Check rate limits
        if not api_optimizer.can_make_request():
            logger.warning(f"Rate limit reached, skipping {operation}")
            return None
        
        # Make API call
        api_optimizer.record_request()
        
        default_params = {
            'max_tokens': 2000,
            'temperature': 0.1
        }
        if params:
            default_params.update(params)
        
        result = call_deepseek_api(
            [{"role": "user", "content": prompt}],
            **default_params
        )
        
        # Cache successful result
        if result:
            api_optimizer.cache_result(cache_key, result)
        
        return result
        
    except Exception as e:
        logger.error(f"Optimized DeepSeek call failed for {operation}: {str(e)}")
        return None

# Enhanced DeepSeek API configuration
# Note: EasyOCR doesn't support Malayalam directly, so we'll use English only for EasyOCR
# and fall back to Tesseract for Malayalam
easyocr_reader = easyocr.Reader(['en'])  # English only for EasyOCR
translator = Translator()

# Configure Tesseract for Malayalam support (Tesseract supports Malayalam as 'mal')
pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

class DocumentChunkProcessor:
    """Intelligent document chunking system for processing entire documents"""
    
    def __init__(self, chunk_size=3000, overlap=500):
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.section_markers = [
            'executive summary', 'introduction', 'conclusion', 'summary', 'recommendation',
            'findings', 'decision', 'background', 'objective', 'purpose', 'scope',
            'methodology', 'analysis', 'result', 'discussion', 'implementation',
            'chapter', 'section', 'part', 'appendix'
        ]
    
    def parse_document_structure(self, text: str) -> Dict:
        """Parse document structure to identify key sections"""
        lines = text.split('\n')
        structure = {
            'sections': [],
            'key_sections': [],
            'total_length': len(text),
            'estimated_chunks': (len(text) // self.chunk_size) + 1
        }
        
        current_pos = 0
        section_map = {}
        
        for line_idx, line in enumerate(lines):
            line_clean = line.strip().lower()
            
            # Look for section headers (lines that contain section markers and are short)
            if line_clean and len(line_clean) < 100:
                for marker in self.section_markers:
                    if marker in line_clean:
                        section_info = {
                            'title': line.strip(),
                            'marker': marker,
                            'position': current_pos,
                            'line_number': line_idx,
                            'importance': self._calculate_section_importance(marker)
                        }
                        structure['sections'].append(section_info)
                        
                        # Mark as key section if important
                        if section_info['importance'] >= 0.7:
                            structure['key_sections'].append(section_info)
                        
                        section_map[current_pos] = section_info
                        break
            
            current_pos += len(line) + 1  # +1 for newline
        
        return structure
    
    def _calculate_section_importance(self, marker: str) -> float:
        """Calculate importance score for section markers"""
        importance_weights = {
            'executive summary': 1.0,
            'summary': 0.9,
            'conclusion': 0.9,
            'recommendation': 0.8,
            'findings': 0.8,
            'decision': 0.8,
            'introduction': 0.7,
            'objective': 0.6,
            'purpose': 0.6,
            'result': 0.7,
            'analysis': 0.6,
            'discussion': 0.5,
            'background': 0.4,
            'methodology': 0.4,
            'implementation': 0.5,
            'scope': 0.3
        }
        return importance_weights.get(marker, 0.3)
    
    def create_intelligent_chunks(self, text: str) -> List[Dict]:
        """Create intelligent chunks with context preservation"""
        if len(text) <= self.chunk_size:
            return [{
                'text': text,
                'chunk_id': 0,
                'start_pos': 0,
                'end_pos': len(text),
                'importance': 1.0,
                'contains_sections': [],
                'is_complete': True
            }]
        
        # Parse document structure
        structure = self.parse_document_structure(text)
        chunks = []
        
        current_pos = 0
        chunk_id = 0
        
        while current_pos < len(text):
            # Calculate chunk end position
            chunk_end = min(current_pos + self.chunk_size, len(text))
            
            # Adjust chunk boundaries to avoid breaking sentences
            if chunk_end < len(text):
                # Look for sentence endings near the boundary
                boundary_text = text[chunk_end-100:chunk_end+100]
                sentence_ends = ['.', '!', '?', '\n\n']
                
                best_break = None
                for end_char in sentence_ends:
                    pos = boundary_text.rfind(end_char)
                    if pos != -1:
                        actual_pos = chunk_end - 100 + pos + 1
                        if chunk_end - 200 < actual_pos < chunk_end + 100:
                            best_break = actual_pos
                            break
                
                if best_break:
                    chunk_end = best_break
            
            # Extract chunk text with overlap
            chunk_start = max(0, current_pos - (self.overlap if chunk_id > 0 else 0))
            chunk_text = text[chunk_start:chunk_end].strip()
            
            # Calculate chunk importance based on contained sections
            chunk_importance = self._calculate_chunk_importance(
                chunk_start, chunk_end, structure['sections']
            )
            
            # Identify sections in this chunk
            contained_sections = []
            for section in structure['sections']:
                if chunk_start <= section['position'] <= chunk_end:
                    contained_sections.append(section['title'])
            
            chunk_info = {
                'text': chunk_text,
                'chunk_id': chunk_id,
                'start_pos': chunk_start,
                'end_pos': chunk_end,
                'importance': chunk_importance,
                'contains_sections': contained_sections,
                'is_complete': chunk_end >= len(text),
                'length': len(chunk_text)
            }
            
            chunks.append(chunk_info)
            current_pos = chunk_end
            chunk_id += 1
        
        return chunks
    
    def _calculate_chunk_importance(self, start_pos: int, end_pos: int, sections: List[Dict]) -> float:
        """Calculate importance score for a chunk based on contained sections"""
        max_importance = 0.3  # Base importance
        
        for section in sections:
            if start_pos <= section['position'] <= end_pos:
                max_importance = max(max_importance, section['importance'])
        
        return max_importance
    
    def prioritize_chunks_for_processing(self, chunks: List[Dict], max_chunks: int = None) -> List[Dict]:
        """Prioritize chunks for processing based on importance"""
        # Sort by importance (descending) and position (ascending for tie-breaking)
        sorted_chunks = sorted(chunks, key=lambda x: (-x['importance'], x['start_pos']))
        
        if max_chunks:
            return sorted_chunks[:max_chunks]
        
        return sorted_chunks
    
    def create_progressive_summary(self, chunks: List[Dict], existing_summary: str = "") -> str:
        """Create a progressive summary by combining chunk summaries"""
        if not chunks:
            return existing_summary
        
        # Combine chunk information
        chunk_summaries = []
        total_sections = set()
        
        for chunk in chunks:
            if chunk['contains_sections']:
                chunk_summaries.append(f"Sections: {', '.join(chunk['contains_sections'])}")
                total_sections.update(chunk['contains_sections'])
        
        structure_info = f"Document contains {len(chunks)} sections including: {', '.join(list(total_sections)[:10])}"
        
        if existing_summary:
            return f"{existing_summary}\n\nAdditional context: {structure_info}"
        else:
            return structure_info

# Initialize document chunk processor
chunk_processor = DocumentChunkProcessor()

class MultilingualDocumentProcessor:
    def __init__(self):
        self.extracted_images = {}  # Store extracted images by document ID
        self.image_metadata = {}    # Store image metadata
        
    def extract_text_and_images_from_pdf(self, pdf_path: str, document_id: str) -> Dict:
        """Extract both text and images from PDF with multilingual support"""
        try:
            # Ensure document_id is string for consistent storage
            document_id = str(document_id)
            
            result = {
                'text_content': '',
                'images': [],
                'text_blocks': [],
                'languages_detected': set(),
                'page_count': 0
            }
            
            # Method 1: Use PyMuPDF for comprehensive extraction
            doc = fitz.open(pdf_path)
            result['page_count'] = len(doc)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract text
                page_text = page.get_text()
                if page_text.strip():
                    result['text_content'] += f"\n--- Page {page_num + 1} ---\n"
                    result['text_content'] += page_text
                    
                    # Detect language of this page
                    try:
                        lang = detect(page_text)
                        result['languages_detected'].add(lang)
                    except:
                        pass
                
                # Extract images with enhanced filtering and quality improvement
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        # Get image data
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        # Skip very small images (likely decorative elements)
                        if pix.width < 50 or pix.height < 50:
                            pix = None
                            continue
                            
                        # Skip images with low quality (too few pixels)
                        if pix.width * pix.height < 10000:  # Less than 100x100 equivalent
                            pix = None
                            continue
                        
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            # Convert to higher quality PNG
                            img_data = pix.tobytes("png")
                            pil_image = Image.open(io.BytesIO(img_data))
                            
                            # Enhance image quality for better OCR
                            enhanced_image = self.enhance_image_for_ocr(pil_image)
                            
                            # Save enhanced image as bytes
                            img_buffer = io.BytesIO()
                            enhanced_image.save(img_buffer, format='PNG', quality=95, optimize=True)
                            enhanced_img_data = img_buffer.getvalue()
                            
                            # Generate unique image ID
                            image_id = f"{document_id}_page{page_num + 1}_img{img_index + 1}"
                            
                            # Extract text from enhanced image using multilingual OCR
                            image_text = self.extract_text_from_image(enhanced_image)
                            
                            # Calculate image quality metrics
                            quality_score = self.calculate_image_quality(pil_image)
                            
                            # Get image coordinates on page
                            img_rect = page.get_image_rects(xref)
                            bbox = img_rect[0] if img_rect else None
                            
                            # Store image with enhanced metadata
                            image_info = {
                                'id': image_id,
                                'page': page_num + 1,
                                'data': base64.b64encode(enhanced_img_data).decode('utf-8'),
                                'format': 'png',
                                'width': enhanced_image.width,
                                'height': enhanced_image.height,
                                'quality_score': quality_score,
                                'text_content': image_text['text'],
                                'languages': image_text['languages'],
                                'confidence': image_text['confidence'],
                                'bbox': [bbox.x0, bbox.y0, bbox.x1, bbox.y1] if bbox else None,
                                'file_size': len(enhanced_img_data),
                                'original_size': len(img_data),
                                'extraction_method': 'pymupdf_enhanced'
                            }
                            
                            result['images'].append(image_info)
                            
                            # Add image text to main content with better formatting
                            if image_text['text'].strip():
                                result['text_content'] += f"\n\n[Image {image_id}]\n"
                                result['text_content'] += f"Location: Page {page_num + 1}\n"
                                result['text_content'] += f"Text: {image_text['text']}\n"
                                result['text_content'] += f"Languages: {', '.join(image_text['languages'])}\n"
                                result['languages_detected'].update(image_text['languages'])
                        
                        pix = None
                    except Exception as e:
                        logger.error(f"Error extracting image {img_index} from page {page_num}: {str(e)}")
            
            doc.close()
            
            # Method 2: Use pdfplumber for better text structure
            try:
                with pdfplumber.open(pdf_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        # Extract structured text blocks
                        words = page.extract_words()
                        for word in words:
                            result['text_blocks'].append({
                                'text': word['text'],
                                'page': page_num + 1,
                                'bbox': [word['x0'], word['top'], word['x1'], word['bottom']],
                                'font': word.get('fontname', ''),
                                'size': word.get('size', 0)
                            })
            except Exception as e:
                logger.warning(f"pdfplumber extraction failed: {str(e)}")
            
            # Store extracted images in memory and database
            self.extracted_images[document_id] = result['images']
            logger.info(f"Stored {len(result['images'])} images for document {document_id} in memory cache")
            
            # Also save to database for persistence
            if result['images']:
                self.save_images_to_database(document_id, result['images'])
            
            logger.info(f"Total documents with images in cache: {len(self.extracted_images)}")
            
            result['languages_detected'] = list(result['languages_detected'])
            return result
            
        except Exception as e:
            logger.error(f"PDF extraction error: {str(e)}")
            return {
                'text_content': '',
                'images': [],
                'text_blocks': [],
                'languages_detected': ['en'],
                'page_count': 0,
                'error': str(e)
            }
    
    def enhance_image_for_ocr(self, image):
        """Enhance image quality for better OCR results"""
        try:
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Resize if too small (improve readability)
            min_size = 300
            if image.width < min_size or image.height < min_size:
                scale = max(min_size / image.width, min_size / image.height)
                new_width = int(image.width * scale)
                new_height = int(image.height * scale)
                image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # Enhance contrast and sharpness
            from PIL import ImageEnhance
            
            # Enhance contrast
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(1.2)
            
            # Enhance sharpness
            enhancer = ImageEnhance.Sharpness(image)
            image = enhancer.enhance(1.1)
            
            return image
        except Exception as e:
            print(f"Error enhancing image: {e}")
            return image

    def calculate_image_quality(self, image):
        """Calculate a quality score for the image"""
        try:
            # Basic quality metrics
            width, height = image.size
            pixel_count = width * height
            
            # Calculate aspect ratio reasonableness
            aspect_ratio = max(width, height) / min(width, height)
            aspect_score = max(0, 1 - (aspect_ratio - 1) / 10)  # Penalty for extreme ratios
            
            # Size score (larger images generally better for OCR)
            size_score = min(1.0, pixel_count / 100000)  # Normalize to 100k pixels
            
            # Combine scores
            quality_score = (aspect_score * 0.3 + size_score * 0.7)
            
            return round(quality_score, 3)
        except Exception as e:
            print(f"Error calculating image quality: {e}")
            return 0.5
    
    def extract_text_from_image(self, image: Image.Image) -> Dict:
        """Extract text from image using both Tesseract and EasyOCR with proper Malayalam support"""
        result = {
            'text': '',
            'languages': [],
            'confidence': 0
        }
        
        try:
            # Convert PIL Image to numpy array for OpenCV
            image_array = np.array(image)
            
            # Method 1: Try Tesseract for Malayalam first
            tesseract_text_ml = ''
            tesseract_text_en = ''
            
            try:
                # Try Malayalam with Tesseract (supports 'mal' language code)
                tesseract_text_ml = pytesseract.image_to_string(image, lang='mal')
                # Also try English
                tesseract_text_en = pytesseract.image_to_string(image, lang='eng')
            except Exception as e:
                logger.warning(f"Tesseract failed: {str(e)}")
            
            # Method 2: EasyOCR for English (more accurate than Tesseract for English)
            easyocr_text = ''
            easyocr_confidence = 0
            
            try:
                easyocr_results = easyocr_reader.readtext(image_array)
                if easyocr_results:
                    easyocr_text = ' '.join([item[1] for item in easyocr_results if item[2] > 0.5])
                    easyocr_confidence = np.mean([item[2] for item in easyocr_results if item[2] > 0.5])
            except Exception as e:
                logger.warning(f"EasyOCR failed: {str(e)}")
            
            # Determine the best result
            candidates = []
            
            # Add Malayalam result if substantial
            if tesseract_text_ml.strip() and len(tesseract_text_ml.strip()) > 5:
                candidates.append({
                    'text': tesseract_text_ml,
                    'languages': ['ml'],
                    'confidence': 0.7,
                    'source': 'tesseract_ml'
                })
            
            # Add English results
            if easyocr_text.strip():
                candidates.append({
                    'text': easyocr_text,
                    'languages': ['en'],
                    'confidence': easyocr_confidence,
                    'source': 'easyocr_en'
                })
            
            if tesseract_text_en.strip():
                candidates.append({
                    'text': tesseract_text_en,
                    'languages': ['en'],
                    'confidence': 0.6,
                    'source': 'tesseract_en'
                })
            
            # Choose the best candidate based on length and confidence
            if candidates:
                best_candidate = max(candidates, key=lambda x: len(x['text'].strip()) * x['confidence'])
                result.update(best_candidate)
                
                # Enhanced OCR post-processing using DeepSeek API
                try:
                    ocr_enhancement = deepseek_ocr_enhancement(
                        result['text'], 
                        f"Image from {best_candidate['source']}"
                    )
                    
                    if ocr_enhancement['confidence'] > result['confidence']:
                        logger.info(f"OCR enhanced: {len(ocr_enhancement['corrections_made'])} corrections made")
                        result['text'] = ocr_enhancement['enhanced_text']
                        result['confidence'] = min(result['confidence'] + 0.1, 1.0)  # Boost confidence
                        result['ocr_corrections'] = ocr_enhancement['corrections_made']
                        result['detected_entities'] = ocr_enhancement['entities']
                        
                except Exception as e:
                    logger.warning(f"OCR enhancement failed: {str(e)}")
                    result['ocr_corrections'] = []
                    result['detected_entities'] = []
                
                # Additional language detection
                try:
                    detected_lang = detect(result['text'])
                    if detected_lang not in result['languages']:
                        result['languages'].append(detected_lang)
                except:
                    pass
                
                # Check for Malayalam characters
                if any('\u0d00' <= char <= '\u0d7f' for char in result['text']):
                    if 'ml' not in result['languages']:
                        result['languages'].append('ml')
            
        except Exception as e:
            logger.error(f"Image text extraction error: {str(e)}")
            result['error'] = str(e)
        
        return result
    
    def translate_text(self, text: str, target_lang: str = 'en') -> str:
        """Translate text to target language using enhanced DeepSeek translation"""
        try:
            if not text.strip():
                return text
            
            # Use enhanced DeepSeek translation
            translation_result = deepseek_enhanced_translation(text, target_lang)
            return translation_result['translated_text']
            
        except Exception as e:
            logger.error(f"Enhanced translation error: {str(e)}")
            # Fallback to Google Translate
            try:
                detected = translator.detect(text)
                if detected.lang == target_lang:
                    return text
                translated = translator.translate(text, dest=target_lang)
                return translated.text
            except:
                return text  # Return original if all translation fails
    
    def get_document_images(self, document_id: str) -> List[Dict]:
        """Get all images for a document with UUID mapping support"""
        # Ensure document_id is string for consistent lookup
        document_id = str(document_id)
        
        # First try to get from database
        try:
            images_from_db = self.get_images_from_database(document_id)
            if images_from_db:
                logger.info(f"Getting images for document {document_id}: Found {len(images_from_db)} images from database")
                return images_from_db
        except Exception as e:
            logger.warning(f"Failed to get images from database for document {document_id}: {str(e)}")
        
        # If not found and document_id is numeric, try to find UUID mapping
        if document_id.isdigit():
            uuid_mapping = self.find_uuid_for_numeric_id(document_id)
            if uuid_mapping:
                logger.info(f"Found UUID mapping for numeric ID {document_id}: {uuid_mapping}")
                try:
                    images_from_db = self.get_images_from_database(uuid_mapping)
                    if images_from_db:
                        return images_from_db
                except Exception as e:
                    logger.warning(f"Failed to get images from database with UUID {uuid_mapping}: {str(e)}")
                
                # Also try memory cache with UUID
                if uuid_mapping in self.extracted_images:
                    images = self.extracted_images[uuid_mapping]
                    logger.info(f"Found {len(images)} images in memory cache with UUID {uuid_mapping}")
                    return images

        # Fallback to memory cache with original ID
        images = self.extracted_images.get(document_id, [])
        logger.info(f"Getting images for document {document_id}: Found {len(images)} images in memory cache")
        logger.info(f"Available document IDs in cache: {list(self.extracted_images.keys())}")
        if len(images) > 0:
            logger.info(f"Sample image keys: {list(images[0].keys()) if images else 'No images'}")
        return images

    def find_uuid_for_numeric_id(self, numeric_id: str) -> str:
        """Find UUID document ID for a given numeric ID"""
        try:
            # Check in-memory cache first for any UUID that has images
            for uuid_key in self.extracted_images.keys():
                if self.extracted_images[uuid_key]:  # Has images
                    # For now, we'll use a simple mapping strategy
                    # In production, you'd have a proper mapping table
                    logger.info(f"Mapping numeric ID {numeric_id} to UUID {uuid_key}")
                    return uuid_key
            return None
        except Exception as e:
            logger.error(f"Error finding UUID mapping: {str(e)}")
            return None
    
    def find_relevant_images(self, document_id: str, query: str, language: str = 'en') -> List[Dict]:
        """Find images relevant to a query with enhanced matching"""
        # Ensure document_id is string for consistent lookup
        document_id = str(document_id)
        images = self.get_document_images(document_id)
        logger.info(f"Finding relevant images for document {document_id}: Found {len(images)} total images")
        
        if not images:
            logger.warning(f"No images found for document {document_id}")
            return []
        
        relevant_images = []
        
        # Translate query to both languages for better matching
        query_en = query if language == 'en' else self.translate_text(query, 'en')
        query_ml = query if language == 'ml' else self.translate_text(query, 'ml')
        
        # Create expanded query words including synonyms
        query_words = set()
        for q in [query.lower(), query_en.lower(), query_ml.lower()]:
            query_words.update(q.split())
        
        logger.info(f"Query words for image search: {query_words}")
        
        for image in images:
            image_text = image.get('text_content', '').lower()
            image_words = set(image_text.split())
            
            # Calculate multiple relevance factors
            factors = []
            
            # 1. Direct word matching
            common_words = query_words.intersection(image_words)
            word_score = len(common_words) / max(len(query_words), 1) if query_words else 0
            factors.append(word_score * 0.6)
            
            # 2. Substring matching for partial matches
            substring_score = 0
            for q_word in query_words:
                if len(q_word) > 3:  # Only for meaningful words
                    for i_word in image_words:
                        if q_word in i_word or i_word in q_word:
                            substring_score += 0.5
            substring_score = min(substring_score / max(len(query_words), 1), 1.0)
            factors.append(substring_score * 0.3)
            
            # 3. Quality bonus for high-quality images
            quality_score = image.get('quality_score', 0.5)
            factors.append(quality_score * 0.1)
            
            # Calculate final relevance score
            relevance_score = sum(factors)
            
            # More lenient threshold for including images
            if relevance_score > 0.05 or len(image_text.strip()) > 10:
                image_copy = image.copy()
                image_copy['relevance_score'] = round(relevance_score, 3)
                relevant_images.append(image_copy)
        
        # Sort by relevance score
        relevant_images.sort(key=lambda x: x['relevance_score'], reverse=True)
        
        logger.info(f"Found {len(relevant_images)} relevant images for query '{query}'")
        for img in relevant_images[:3]:  # Log top 3
            logger.info(f"  - Image {img.get('id', 'unknown')} (Page {img.get('page', 'N/A')}): Score {img['relevance_score']}")
        
        # Return more images if available, prioritizing quality
        return relevant_images[:8]  # Increased from 5 to 8
    
    def save_images_to_database(self, document_id: str, images: List[Dict]):
        """Save extracted images to database with support for both numeric and UUID document IDs"""
        try:
            # Only save to database if document_id is numeric (backend database format)
            if not document_id.isdigit():
                logger.info(f"Skipping database save for UUID document_id {document_id} - storing in memory cache only")
                return
            
            conn = get_db_connection()
            cursor = conn.cursor()
            
            for image in images:
                # Insert image into document_images table
                cursor.execute("""
                    INSERT INTO document_images 
                    (document_id, image_id, page_number, image_data, image_format, 
                     text_content, languages_detected, bbox, extraction_confidence)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
                    ON CONFLICT (document_id, image_id) DO UPDATE SET
                        image_data = EXCLUDED.image_data,
                        text_content = EXCLUDED.text_content,
                        extraction_confidence = EXCLUDED.extraction_confidence
                """, (
                    int(document_id),
                    image['id'],
                    image.get('page', 1),
                    image['data'],
                    image.get('format', 'png'),
                    image.get('text_content', ''),
                    json.dumps(image.get('languages', [])),
                    json.dumps(image.get('bbox', [])),
                    image.get('confidence', 0.0)
                ))
            
            conn.commit()
            cursor.close()
            conn.close()
            logger.info(f"Saved {len(images)} images to database for document {document_id}")
            
        except Exception as e:
            logger.error(f"Failed to save images to database: {str(e)}")
    
    def get_images_from_database(self, document_id: str) -> List[Dict]:
        """Retrieve images from database with support for both numeric and UUID document IDs"""
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            
            # Try to query with the document_id as provided
            try:
                # If document_id is numeric, cast to int for database query
                if document_id.isdigit():
                    cursor.execute("""
                        SELECT image_id, page_number, image_data, image_format, 
                               text_content, languages_detected, bbox, extraction_confidence
                        FROM document_images 
                        WHERE document_id = %s
                        ORDER BY page_number, image_id
                    """, (int(document_id),))
                else:
                    # For UUID document_ids, we need to handle them differently
                    # Since our current database schema expects numeric IDs, 
                    # UUID documents won't be in the database yet
                    logger.info(f"UUID document_id {document_id} not supported in current database schema")
                    cursor.close()
                    conn.close()
                    return []
                
                rows = cursor.fetchall()
                cursor.close()
                conn.close()
                
                if not rows:
                    logger.info(f"No images found in database for document {document_id}")
                    return []
                
                images = []
                for row in rows:
                    image = {
                        'id': row['image_id'],
                        'page': row['page_number'],
                        'data': row['image_data'],
                        'format': row['image_format'],
                        'text_content': row['text_content'] or '',
                        'languages': json.loads(row['languages_detected'] or '[]'),
                        'bbox': json.loads(row['bbox'] or '[]'),
                        'confidence': row['extraction_confidence'] or 0.0
                    }
                    images.append(image)
                
                logger.info(f"Retrieved {len(images)} images from database for document {document_id}")
                return images
                
            except Exception as query_error:
                cursor.close()
                conn.close()
                logger.error(f"Database query failed for document {document_id}: {str(query_error)}")
                return []
            
        except Exception as e:
            logger.error(f"Failed to retrieve images from database: {str(e)}")
            return []
            
            images = []
            for row in rows:
                images.append({
                    'id': row['image_id'],
                    'page': row['page_number'],
                    'data': row['image_data'],
                    'format': row['image_format'],
                    'text_content': row['text_content'] or '',
                    'languages': json.loads(row['languages_detected']) if row['languages_detected'] else [],
                    'bbox': json.loads(row['bbox']) if row['bbox'] else [],
                    'confidence': row['extraction_confidence'] or 0.0
                })
            
            return images
            
        except Exception as e:
            logger.error(f"Failed to get images from database: {str(e)}")
            return []
    
    def load_all_images_from_database(self):
        """Load all images from database into memory cache on startup"""
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            
            cursor.execute("""
                SELECT DISTINCT document_id FROM document_images
            """)
            
            document_ids = [str(row['document_id']) for row in cursor.fetchall()]
            cursor.close()
            conn.close()
            
            for doc_id in document_ids:
                images = self.get_images_from_database(doc_id)
                if images:
                    self.extracted_images[doc_id] = images
                    
            logger.info(f"Loaded images for {len(document_ids)} documents from database into cache")
            
        except Exception as e:
            logger.error(f"Failed to load images from database: {str(e)}")

# Initialize the multilingual processor
doc_processor = MultilingualDocumentProcessor()

# Database connection
def get_db_connection():
    return psycopg2.connect(
        os.getenv('DATABASE_URL'),
        cursor_factory=RealDictCursor
    )

# Load existing images from database into cache on startup
try:
    doc_processor.load_all_images_from_database()
    logger.info("Successfully loaded existing images from database into cache")
except Exception as e:
    logger.error(f"Failed to load images from database on startup: {str(e)}")

class DocumentProcessor:
    def __init__(self):
        self.role_keywords = {
            'LEADERSHIP': [
                'board', 'director', 'executive', 'strategy', 'policy', 'decision',
                'approval', 'management', 'leadership', 'governance', 'vision',
                'mission', 'objectives', 'planning', 'budget approval', 'oversight'
            ],
            'HR': [
                'employee', 'staff', 'recruitment', 'hiring', 'payroll', 'salary',
                'benefits', 'training', 'performance', 'appraisal', 'leave',
                'attendance', 'resignation', 'termination', 'grievance', 'policy'
            ],
            'FINANCE': [
                'budget', 'finance', 'accounting', 'audit', 'expense', 'revenue',
                'cost', 'invoice', 'payment', 'procurement', 'vendor', 'contract',
                'financial', 'expenditure', 'income', 'profit', 'loss', 'balance'
            ],
            'ENGINEER': [
                'technical', 'engineering', 'design', 'specification', 'maintenance',
                'construction', 'project', 'infrastructure', 'system', 'equipment',
                'safety', 'quality', 'testing', 'inspection', 'drawing', 'blueprint'
            ]
        }
    
    def extract_text_from_file(self, file_path: str, mime_type: str) -> str:
        """Extract text from various file formats."""
        try:
            if mime_type == 'application/pdf':
                return self._extract_from_pdf(file_path)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return self._extract_from_docx(file_path)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return self._extract_from_pptx(file_path)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                return self._extract_from_xlsx(file_path)
            elif mime_type == 'application/vnd.ms-excel':
                return self._extract_from_xls(file_path)
            elif mime_type == 'application/vnd.ms-powerpoint':
                return self._extract_from_ppt(file_path)
            elif mime_type == 'application/msword':
                return self._extract_from_doc(file_path)
            elif mime_type == 'text/plain':
                return self._extract_from_txt(file_path)
            elif mime_type == 'text/csv':
                return self._extract_from_csv(file_path)
            elif mime_type == 'application/rtf':
                return self._extract_from_rtf(file_path)
            elif mime_type.startswith('image/'):
                return self._extract_from_image(file_path)
            else:
                return ""
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            return ""
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            logger.error(f"PDF extraction error: {str(e)}")
        return text.strip()
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
        except Exception as e:
            logger.error(f"DOCX extraction error: {str(e)}")
            return ""
    
    def _extract_from_image(self, file_path: str) -> str:
        """Extract text from image using OCR."""
        try:
            # Configure Tesseract for English and Malayalam
            custom_config = r'--oem 3 --psm 6 -l eng+mal'
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, config=custom_config)
            return text.strip()
        except Exception as e:
            logger.error(f"OCR extraction error: {str(e)}")
            return ""

    def _extract_from_txt(self, file_path: str) -> str:
        """Extract text from plain text file."""
        try:
            # Detect encoding first
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            
            # Read with detected encoding
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read().strip()
        except Exception as e:
            logger.error(f"TXT extraction error: {str(e)}")
            return ""

    def _extract_from_csv(self, file_path: str) -> str:
        """Extract text from CSV file."""
        try:
            # Detect encoding first
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            
            # Read CSV and convert to text
            df = pd.read_csv(file_path, encoding=encoding)
            # Convert dataframe to readable text format
            text_lines = []
            text_lines.append("CSV Data:")
            text_lines.append("Columns: " + ", ".join(df.columns.tolist()))
            text_lines.append("")
            
            for index, row in df.iterrows():
                row_text = " | ".join([f"{col}: {str(val)}" for col, val in row.items()])
                text_lines.append(row_text)
                
            return "\n".join(text_lines)
        except Exception as e:
            logger.error(f"CSV extraction error: {str(e)}")
            return ""

    def _extract_from_xlsx(self, file_path: str) -> str:
        """Extract text from Excel XLSX file."""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_lines = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_lines.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text_lines.append(row_text)
                text_lines.append("")
                
            return "\n".join(text_lines)
        except Exception as e:
            logger.error(f"XLSX extraction error: {str(e)}")
            return ""

    def _extract_from_xls(self, file_path: str) -> str:
        """Extract text from Excel XLS file."""
        try:
            workbook = xlrd.open_workbook(file_path)
            text_lines = []
            
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                text_lines.append(f"Sheet: {sheet.name}")
                
                for row_idx in range(sheet.nrows):
                    row_values = []
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        row_values.append(str(cell_value) if cell_value else "")
                    
                    row_text = " | ".join(row_values)
                    if row_text.strip():
                        text_lines.append(row_text)
                text_lines.append("")
                
            return "\n".join(text_lines)
        except Exception as e:
            logger.error(f"XLS extraction error: {str(e)}")
            return ""

    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint PPTX file."""
        try:
            presentation = Presentation(file_path)
            text_lines = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_lines.append(f"Slide {slide_num}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_lines.append(shape.text)
                
                text_lines.append("")
                
            return "\n".join(text_lines)
        except Exception as e:
            logger.error(f"PPTX extraction error: {str(e)}")
            return ""

    def _extract_from_ppt(self, file_path: str) -> str:
        """Extract text from PowerPoint PPT file."""
        try:
            # For older PPT files, we'll try to convert using python-pptx
            # This might not work for all PPT files, but it's a fallback
            logger.warning(f"PPT format not fully supported, attempting basic extraction: {file_path}")
            return "PPT file detected. Please convert to PPTX format for better text extraction."
        except Exception as e:
            logger.error(f"PPT extraction error: {str(e)}")
            return ""

    def _extract_from_doc(self, file_path: str) -> str:
        """Extract text from Word DOC file."""
        try:
            # For older DOC files, we'll provide a message
            logger.warning(f"DOC format not fully supported, basic extraction only: {file_path}")
            return "DOC file detected. Please convert to DOCX format for better text extraction."
        except Exception as e:
            logger.error(f"DOC extraction error: {str(e)}")
            return ""

    def _extract_from_rtf(self, file_path: str) -> str:
        """Extract text from RTF file."""
        try:
            # Basic RTF extraction - remove RTF formatting
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                
            # Simple RTF tag removal (basic approach)
            import re
            # Remove RTF control words
            content = re.sub(r'\\[a-z]+\d*\s*', '', content)
            # Remove curly braces
            content = re.sub(r'[{}]', '', content)
            # Clean up extra whitespace
            content = re.sub(r'\s+', ' ', content)
            
            return content.strip()
        except Exception as e:
            logger.error(f"RTF extraction error: {str(e)}")
            return ""
    
    def detect_language(self, text: str) -> str:
        """Detect the primary language of the text using enhanced DeepSeek detection."""
        try:
            # Use enhanced DeepSeek language detection
            lang_info = deepseek_enhanced_language_detection(text)
            return lang_info['language']
        except:
            # Fallback to local detection
            try:
                clean_text = re.sub(r'[^\w\s]', ' ', text)
                if len(clean_text.strip()) < 10:
                    return 'en'
                detected = detect(clean_text)
                return detected
            except:
                return 'en'
    
    def generate_ai_analysis(self, text: str) -> Dict:
        """Generate comprehensive AI analysis of entire document using intelligent chunking."""
        try:
            # Detect language of full document
            primary_language = self.detect_language(text)
            logger.info(f"Starting enhanced AI analysis for document with {len(text)} characters")
            
            # Create intelligent chunks for the entire document
            chunks = chunk_processor.create_intelligent_chunks(text)
            logger.info(f"Created {len(chunks)} chunks for processing")
            
            # Prioritize chunks by importance (process most important first)
            prioritized_chunks = chunk_processor.prioritize_chunks_for_processing(chunks, max_chunks=8)
            
            # Process chunks progressively
            chunk_analyses = []
            combined_insights = {
                'key_topics': set(),
                'entities': set(),
                'document_types': [],
                'sensitivity_indicators': [],
                'role_indicators': []
            }
            
            for i, chunk in enumerate(prioritized_chunks):
                logger.info(f"Processing chunk {i+1}/{len(prioritized_chunks)} (importance: {chunk['importance']:.2f})")
                
                # Create chunk-specific prompt
                chunk_prompt = f"""
                Analyze this section of a KMRL document and extract key information:
                
                Document Section (Chunk {i+1}/{len(prioritized_chunks)}):
                Sections contained: {', '.join(chunk['contains_sections']) if chunk['contains_sections'] else 'Content section'}
                Importance level: {chunk['importance']:.2f}
                
                Text:
                {chunk['text'][:3500]}  # Use larger chunk for analysis
                
                Extract and return JSON (use only valid KMRL roles: LEADERSHIP, HR, FINANCE, ENGINEER, ADMIN):
                {{
                    "section_summary": "Summary of this section's content",
                    "key_topics": ["topic1", "topic2"],
                    "entities": ["entity1", "entity2"],
                    "document_type_indicators": ["contract", "policy"],
                    "sensitivity_indicators": ["confidential", "internal"],
                    "role_relevance": ["LEADERSHIP", "HR", "FINANCE", "ENGINEER", "ADMIN"],
                    "important_details": ["detail1", "detail2"]
                }}
                
                Focus on extracting the most important information from this section.
                """
                
                try:
                    chunk_response = optimized_deepseek_call(
                        f"chunk_analysis_{i}",
                        chunk['text'][:500],  # Cache key based on chunk content
                        chunk_prompt,
                        {'max_tokens': 1500, 'temperature': 0.2}
                    )
                    
                    if chunk_response:
                        # Parse chunk analysis
                        chunk_response = chunk_response.strip()
                        if chunk_response.startswith('```json'):
                            chunk_response = chunk_response[7:]
                        if chunk_response.endswith('```'):
                            chunk_response = chunk_response[:-3]
                        
                        chunk_data = json.loads(chunk_response)
                        chunk_analyses.append({
                            'chunk_id': chunk['chunk_id'],
                            'importance': chunk['importance'],
                            'sections': chunk['contains_sections'],
                            'analysis': chunk_data
                        })
                        
                        # Aggregate insights
                        if 'key_topics' in chunk_data:
                            combined_insights['key_topics'].update(chunk_data['key_topics'])
                        if 'entities' in chunk_data:
                            combined_insights['entities'].update(chunk_data['entities'])
                        if 'document_type_indicators' in chunk_data:
                            combined_insights['document_types'].extend(chunk_data['document_type_indicators'])
                        if 'sensitivity_indicators' in chunk_data:
                            combined_insights['sensitivity_indicators'].extend(chunk_data['sensitivity_indicators'])
                        if 'role_relevance' in chunk_data:
                            combined_insights['role_indicators'].extend(chunk_data['role_relevance'])
                        
                    else:
                        logger.warning(f"No response for chunk {i+1}")
                        
                except Exception as e:
                    logger.warning(f"Failed to process chunk {i+1}: {str(e)}")
                    continue
            
            # Create comprehensive summary from all chunks
            logger.info("Creating comprehensive summary from chunk analyses")
            
            # Prepare comprehensive analysis prompt
            chunk_summaries = []
            for analysis in chunk_analyses:
                summary = analysis['analysis'].get('section_summary', '')
                sections = ', '.join(analysis['sections']) if analysis['sections'] else 'Content'
                chunk_summaries.append(f"[{sections}]: {summary}")
            
            comprehensive_prompt = f"""
            Based on analysis of {len(chunks)} sections of this KMRL document, create a comprehensive analysis:
            
            Document Overview:
            - Total length: {len(text)} characters
            - Sections analyzed: {len(chunk_analyses)}
            - Primary language: {primary_language}
            
            Section Summaries:
            {chr(10).join(chunk_summaries[:10])}  # Limit to first 10 summaries
            
            Aggregated Insights:
            - Key topics: {', '.join(list(combined_insights['key_topics'])[:15])}
            - Entities: {', '.join(list(combined_insights['entities'])[:10])}
            - Document type indicators: {', '.join(set(combined_insights['document_types']))}
            - Sensitivity indicators: {', '.join(set(combined_insights['sensitivity_indicators']))}
            
            Create final analysis STRICTLY in JSON format (no prose). IMPORTANT: Always include a numeric confidence between 0 and 1 under recommended_roles.confidence. If unsure, estimate reasonably based on evidence.
            
            CRITICAL: For recommended_roles.roles, you MUST ONLY use these exact values: LEADERSHIP, HR, FINANCE, ENGINEER, ADMIN. Do not create or suggest any other role names.
            
            {{
                "summary_en": "Comprehensive summary covering all analyzed sections",
                "summary_ml": "Malayalam summary (if applicable)",
                "tags": ["tag1", "tag2", "tag3"],
                "primary_language": "{primary_language}",
                "sensitivity_level": "LOW|MEDIUM|HIGH|CONFIDENTIAL",
                "document_type": "contract|policy|report|invoice|memo|other",
                "key_entities": ["entity1", "entity2"],
                "recommended_roles": {{
                    "roles": ["LEADERSHIP", "HR", "FINANCE", "ENGINEER", "ADMIN"],
                    "confidence": 0.85,
                    "reasoning": "Explanation based on full document analysis"
                }},
                "retention_recommendation": {{
                    "days": 2555,
                    "reason": "Legal/business justification"
                }},
                "analysis_metadata": {{
                    "chunks_processed": {len(chunk_analyses)},
                    "total_chunks": {len(chunks)},
                    "coverage_percentage": {round((len(chunk_analyses) / len(chunks)) * 100, 1)},
                    "processing_method": "enhanced_chunking"
                }}
            }}
            
            Ensure the summary represents the ENTIRE document based on all processed sections.
            """
            
            try:
                final_response = optimized_deepseek_call(
                    "comprehensive_analysis",
                    f"{len(text)}_{len(chunk_analyses)}",  # Cache key based on document size and chunks
                    comprehensive_prompt,
                    {'max_tokens': 4000, 'temperature': 0.3}
                )
                
                if final_response:
                    logger.info(f"Received comprehensive analysis: {len(final_response)} characters")
                    
                    # Parse final response
                    final_response = final_response.strip()
                    if final_response.startswith('```json'):
                        final_response = final_response[7:]
                    if final_response.endswith('```'):
                        final_response = final_response[:-3]
                    
                    analysis = json.loads(final_response)
                    # Normalize confidence to ensure it's always present and valid
                    try:
                        rr = analysis.get('recommended_roles') or {}
                        conf = rr.get('confidence')
                        if conf is None:
                            # Derive a conservative default based on coverage
                            coverage = max(min(len(chunk_analyses) / max(len(chunks), 1), 1.0), 0.0)
                            conf = 0.6 + 0.3 * coverage
                        conf = float(conf)
                        if not (conf == conf):  # NaN check
                            conf = 0.75
                        conf = max(0.0, min(conf, 1.0))
                        rr['confidence'] = conf
                        analysis['recommended_roles'] = rr
                    except Exception as _:
                        analysis.setdefault('recommended_roles', {}).setdefault('confidence', 0.75)
                    logger.info("Successfully parsed comprehensive AI response")
                    
                    # Apply heuristic role assignment as backup
                    heuristic_roles = self._apply_role_heuristics(text)
                    
                    # Combine AI and heuristic results
                    final_roles = self._combine_role_predictions(
                        analysis.get('recommended_roles', {}),
                        heuristic_roles
                    )
                    
                    analysis['recommended_roles'] = final_roles
                    logger.info("Enhanced AI analysis completed successfully")
                    return analysis
                    
                else:
                    raise Exception("No response from comprehensive analysis")
                    
            except Exception as e:
                logger.error(f"Final analysis failed: {str(e)}")
                # Fall back to partial analysis if final synthesis fails
                return self._create_analysis_from_chunks(chunk_analyses, text, primary_language)
                
        except Exception as e:
            logger.error(f"Enhanced AI analysis error: {str(e)}")
            logger.error(traceback.format_exc())
            return self._generate_fallback_analysis(text)
    
    def _create_analysis_from_chunks(self, chunk_analyses: List[Dict], text: str, primary_language: str) -> Dict:
        """Create analysis from chunk data when comprehensive analysis fails"""
        logger.info("Creating analysis from individual chunk results")
        
        # Aggregate data from chunks
        all_topics = set()
        all_entities = set()
        all_summaries = []
        
        for chunk_analysis in chunk_analyses:
            data = chunk_analysis['analysis']
            if 'key_topics' in data:
                all_topics.update(data['key_topics'])
            if 'entities' in data:
                all_entities.update(data['entities'])
            if 'section_summary' in data:
                all_summaries.append(data['section_summary'])
        
        # Create combined summary
        combined_summary = f"Document analysis based on {len(chunk_analyses)} sections: " + \
                          " ".join(all_summaries[:5])  # Limit to prevent overflow
        
        # Apply heuristic analysis for other fields
        heuristic_analysis = self._generate_fallback_analysis(text)
        
        # Enhance with chunk data
        heuristic_analysis['summary_en'] = combined_summary[:1000]  # Limit length
        heuristic_analysis['tags'] = list(all_topics)[:8]
        heuristic_analysis['key_entities'] = list(all_entities)[:6]
        heuristic_analysis['primary_language'] = primary_language
        heuristic_analysis['analysis_metadata'] = {
            'chunks_processed': len(chunk_analyses),
            'processing_method': 'chunk_aggregation',
            'fallback_used': True
        }
        
        return heuristic_analysis
    
    def _apply_role_heuristics(self, text: str) -> Dict:
        """Apply keyword-based heuristics for role assignment."""
        text_lower = text.lower()
        role_scores = {}
        
        for role, keywords in self.role_keywords.items():
            score = sum(1 for keyword in keywords if keyword.lower() in text_lower)
            role_scores[role] = score / len(keywords)  # Normalize
        
        # Get roles with significant scores
        threshold = 0.1
        relevant_roles = [role for role, score in role_scores.items() if score > threshold]
        
        if not relevant_roles:
            relevant_roles = ['LEADERSHIP']  # Default fallback
        
        max_score = max(role_scores.values()) if role_scores else 0.1
        confidence = min(max_score * 2, 1.0)  # Scale confidence
        
        return {
            'roles': relevant_roles,
            'confidence': confidence,
            'reasoning': f"Keyword-based analysis identified relevance to: {', '.join(relevant_roles)}"
        }
    
    def _combine_role_predictions(self, ai_prediction: Dict, heuristic_prediction: Dict) -> Dict:
        """Combine AI and heuristic role predictions."""
        ai_roles = set(ai_prediction.get('roles', []))
        heuristic_roles = set(heuristic_prediction.get('roles', []))
        
        # Union of both predictions
        combined_roles = list(ai_roles.union(heuristic_roles))
        
        # Average confidence
        ai_confidence = ai_prediction.get('confidence', 0.5)
        heuristic_confidence = heuristic_prediction.get('confidence', 0.5)
        combined_confidence = (ai_confidence + heuristic_confidence) / 2
        
        return {
            'roles': combined_roles,
            'confidence': combined_confidence,
            'reasoning': f"Combined AI and heuristic analysis. AI suggested: {list(ai_roles)}, Heuristics suggested: {list(heuristic_roles)}"
        }
    
    def _generate_fallback_analysis(self, text: str) -> Dict:
        """Generate enhanced fallback analysis when AI fails."""
        logger.warning("Using enhanced fallback analysis - AI analysis failed")
        heuristic_roles = self._apply_role_heuristics(text)
        
        # Extract key information for better summary
        words = text.split()
        sentences = [s.strip() for s in text.split('.') if s.strip()]
        
        # Enhanced keyword extraction
        common_words = {'the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should', 'this', 'that', 'these', 'those', 'they', 'them', 'their', 'there', 'where', 'when', 'what', 'who', 'why', 'how'}
        significant_words = [word.lower().strip('.,!?;:()[]{}"-') for word in words if word.lower() not in common_words and len(word) > 3 and word.isalpha()]
        
        # Get word frequency and filter out less meaningful words
        from collections import Counter
        word_freq = Counter(significant_words)
        top_keywords = [word for word, count in word_freq.most_common(8) if count > 1]
        
        # Detect document type based on content
        text_lower = text.lower()
        doc_type = 'other'
        
        if any(term in text_lower for term in ['contract', 'agreement', 'terms', 'conditions']):
            doc_type = 'contract'
        elif any(term in text_lower for term in ['policy', 'procedure', 'guideline', 'rules']):
            doc_type = 'policy'
        elif any(term in text_lower for term in ['report', 'analysis', 'findings', 'conclusions']):
            doc_type = 'report'
        elif any(term in text_lower for term in ['invoice', 'bill', 'payment', 'amount', 'total']):
            doc_type = 'invoice'
        elif any(term in text_lower for term in ['memo', 'memorandum', 'notice', 'announcement']):
            doc_type = 'memo'
        
        # Enhanced summary generation
        # Try to find the most informative sentences
        important_sentences = []
        for sentence in sentences[:10]:  # Check first 10 sentences
            if any(keyword in sentence.lower() for keyword in top_keywords[:5]):
                important_sentences.append(sentence)
            if len(important_sentences) >= 3:
                break
        
        if not important_sentences:
            important_sentences = sentences[:3]
        
        # Create structured summary
        if doc_type == 'contract':
            summary_template = f"Contract/Agreement document concerning {', '.join(top_keywords[:3])}. "
        elif doc_type == 'policy':
            summary_template = f"Policy document addressing {', '.join(top_keywords[:3])}. "
        elif doc_type == 'report':
            summary_template = f"Report analyzing {', '.join(top_keywords[:3])}. "
        elif doc_type == 'invoice':
            summary_template = f"Financial document (invoice/bill) related to {', '.join(top_keywords[:3])}. "
        elif doc_type == 'memo':
            summary_template = f"Memo/Notice regarding {', '.join(top_keywords[:3])}. "
        else:
            summary_template = f"Document discussing {', '.join(top_keywords[:3])}. "
        
        key_content = ' '.join(important_sentences)[:300]
        summary = f"{summary_template}{key_content}... Document contains {len(words)} words and covers topics including: {', '.join(top_keywords[:5])}."
        
        # Determine sensitivity based on content keywords and document type
        sensitivity = 'MEDIUM'  # Default
        
        if any(word in text_lower for word in ['confidential', 'restricted', 'classified', 'secret', 'private', 'sensitive']):
            sensitivity = 'CONFIDENTIAL'
        elif any(word in text_lower for word in ['internal', 'limited', 'proprietary', 'draft']) or doc_type in ['contract']:
            sensitivity = 'HIGH'
        elif any(word in text_lower for word in ['public', 'announcement', 'press', 'general', 'open']) or doc_type in ['memo']:
            sensitivity = 'LOW'
        
        # Enhanced tags based on content and document type
        base_tags = [doc_type] if doc_type != 'other' else []
        content_tags = []
        
        # Add domain-specific tags
        if any(term in text_lower for term in ['metro', 'railway', 'train', 'station', 'transport']):
            content_tags.append('transportation')
        if any(term in text_lower for term in ['finance', 'budget', 'cost', 'payment', 'money']):
            content_tags.append('finance')
        if any(term in text_lower for term in ['safety', 'security', 'emergency', 'risk']):
            content_tags.append('safety')
        if any(term in text_lower for term in ['project', 'construction', 'development', 'engineering']):
            content_tags.append('engineering')
        if any(term in text_lower for term in ['staff', 'employee', 'personnel', 'human', 'resource']):
            content_tags.append('hr')
        
        all_tags = base_tags + content_tags + top_keywords[:3]
        
        # Improve confidence for fallback based on content analysis
        base_confidence = heuristic_roles.get('confidence', 0.6)
        
        # Increase confidence if we have clear document type or domain keywords
        if doc_type != 'other':
            base_confidence = min(base_confidence + 0.15, 0.85)
        if content_tags:
            base_confidence = min(base_confidence + 0.1, 0.85)
        
        return {
            'summary_en': summary,
            'summary_ml': "",  # Could be enhanced with translation logic
            'tags': list(set(all_tags))[:8],  # Remove duplicates and limit
            'primary_language': self.detect_language(text),
            'sensitivity_level': sensitivity,
            'document_type': doc_type,
            'key_entities': top_keywords[:6],
            'recommended_roles': {
                'roles': heuristic_roles.get('roles', ['LEADERSHIP']),
                'confidence': base_confidence,
                'reasoning': f"Enhanced heuristic analysis identified {doc_type} document with {sensitivity.lower()} sensitivity. Key topics: {', '.join(top_keywords[:3])}"
            },
            'retention_recommendation': {
                'days': 2555 if doc_type in ['contract', 'policy'] else 1825,
                'reason': f'Standard retention for {doc_type} documents based on KMRL policy'
            }
        }
    
    def generate_embeddings(self, text: str) -> List[float]:
        """Generate enhanced semantic embeddings using DeepSeek API."""
        try:
            # Use enhanced DeepSeek embeddings for better semantic understanding
            return deepseek_enhanced_embeddings(text)
        except Exception as e:
            logger.error(f"Enhanced embedding generation error: {str(e)}")
            # Fallback to local TF-IDF
            return generate_local_embeddings(text)

# Initialize processor
processor = DocumentProcessor()

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint."""
    return jsonify({'status': 'healthy', 'timestamp': datetime.now().isoformat()})

@app.route('/api-status', methods=['GET'])
def api_status():
    """Get DeepSeek API usage statistics and system status"""
    try:
        usage_stats = api_optimizer.get_usage_stats()
        
        # Test API connectivity
        test_response = optimized_deepseek_call(
            "test",
            "test connectivity",
            "Respond with 'OK' to confirm API connectivity.",
            {'max_tokens': 10, 'temperature': 0}
        )
        
        api_healthy = test_response and 'OK' in test_response.upper()
        
        status_info = {
            'deepseek_api': {
                'status': 'healthy' if api_healthy else 'degraded',
                'usage_stats': usage_stats,
                'last_test': datetime.now().isoformat()
            },
            'local_fallbacks': {
                'tesseract_available': os.path.exists(TESSERACT_CMD),
                'easyocr_loaded': easyocr_reader is not None,
                'google_translate_available': translator is not None
            },
            'enhancements_active': {
                'language_detection': True,
                'translation': True,
                'ocr_enhancement': True,
                'semantic_embeddings': True,
                'batch_processing': True
            },
            'system_status': 'optimal',
            'timestamp': datetime.now().isoformat()
        }
        
        return jsonify(status_info)
        
    except Exception as e:
        logger.error(f"API status check failed: {str(e)}")
        return jsonify({
            'status': 'error',
            'error': str(e),
            'timestamp': datetime.now().isoformat()
        }), 500

@app.route('/process-document', methods=['POST'])
def process_document():
    """Process uploaded document with multilingual text and image extraction"""
    try:
        data = request.json
        file_path = data.get('file_path')
        mime_type = data.get('mime_type')
        document_id = data.get('document_id')
        
        if not file_path or not mime_type:
            return jsonify({'error': 'file_path and mime_type are required'}), 400
        
        if not document_id:
            document_id = str(uuid.uuid4())
        else:
            # Ensure document_id is always a string for consistent storage/retrieval
            document_id = str(document_id)
        
        logger.info(f"Processing document {document_id} at path: {file_path}")
        
        # Check if file exists
        if not os.path.exists(file_path):
            return jsonify({'error': f'File not found: {file_path}'}), 404
        
        # Extract text and images based on file type
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            extraction_result = doc_processor.extract_text_and_images_from_pdf(file_path, document_id)
            extracted_text = extraction_result['text_content']
            extracted_images = extraction_result['images']
            languages = extraction_result['languages_detected']
            
            # Ensure images are stored (should already be done in extract_text_and_images_from_pdf)
            logger.info(f"PDF processing completed: {len(extracted_images)} images extracted")
            
        elif file_extension in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
            # Process single image
            with Image.open(file_path) as img:
                image_text = doc_processor.extract_text_from_image(img)
                extracted_text = image_text['text']
                languages = image_text['languages']
                
                # Convert image to base64 for storage
                with open(file_path, 'rb') as img_file:
                    img_data = base64.b64encode(img_file.read()).decode('utf-8')
                
                extracted_images = [{
                    'id': f"{document_id}_main_image",
                    'page': 1,
                    'data': img_data,
                    'format': file_extension[1:],
                    'text_content': extracted_text,
                    'languages': languages
                }]
                
                doc_processor.extracted_images[document_id] = extracted_images
                doc_processor.save_images_to_database(document_id, extracted_images)
                
        elif file_extension in ['.doc', '.docx']:
            # Process Word document
            doc = Document(file_path)
            extracted_text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            
            # Extract images from Word document
            extracted_images = []
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        img = Image.open(io.BytesIO(image_data))
                        image_text = doc_processor.extract_text_from_image(img)
                        
                        image_info = {
                            'id': f"{document_id}_word_img_{len(extracted_images)}",
                            'page': 1,
                            'data': base64.b64encode(image_data).decode('utf-8'),
                            'format': 'png',
                            'text_content': image_text['text'],
                            'languages': image_text['languages']
                        }
                        extracted_images.append(image_info)
                        extracted_text += f"\n[Image Text]: {image_text['text']}"
                    except Exception as e:
                        logger.warning(f"Failed to extract image from Word doc: {str(e)}")
            
            doc_processor.extracted_images[document_id] = extracted_images
            if extracted_images:
                doc_processor.save_images_to_database(document_id, extracted_images)
            languages = [detect(extracted_text)] if extracted_text.strip() else ['en']
            
        elif file_extension in ['.ppt', '.pptx']:
            # Process PowerPoint document
            if file_extension == '.pptx':
                presentation = Presentation(file_path)
                text_lines = []
                extracted_images = []
                
                for slide_num, slide in enumerate(presentation.slides, 1):
                    text_lines.append(f"Slide {slide_num}:")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_lines.append(shape.text)
                        
                        # Extract images from slide
                        if shape.shape_type == 13:  # Picture shape type
                            try:
                                image_data = shape.image.blob
                                img = Image.open(io.BytesIO(image_data))
                                image_text = doc_processor.extract_text_from_image(img)
                                
                                image_info = {
                                    'id': f"{document_id}_pptx_slide_{slide_num}_img_{len(extracted_images)}",
                                    'page': slide_num,
                                    'data': base64.b64encode(image_data).decode('utf-8'),
                                    'format': 'png',
                                    'text_content': image_text['text'],
                                    'languages': image_text['languages']
                                }
                                extracted_images.append(image_info)
                                text_lines.append(f"[Image Text]: {image_text['text']}")
                            except Exception as e:
                                logger.warning(f"Failed to extract image from PPTX slide {slide_num}: {str(e)}")
                    
                    text_lines.append("")
                
                extracted_text = '\n'.join(text_lines)
                doc_processor.extracted_images[document_id] = extracted_images
                if extracted_images:
                    doc_processor.save_images_to_database(document_id, extracted_images)
            else:
                # For older PPT files, only extract text (limited support)
                extracted_text = "PPT file detected. Please convert to PPTX format for better text and image extraction."
                extracted_images = []
            
            languages = [detect(extracted_text)] if extracted_text.strip() else ['en']
            
        elif file_extension in ['.xls', '.xlsx', '.csv']:
            # Process Excel/CSV document
            extracted_images = []  # Excel files typically don't contain extractable images in the same way
            
            try:
                if file_extension == '.csv':
                    # Handle CSV files
                    df = pd.read_csv(file_path)
                    extracted_text = f"CSV Data ({len(df)} rows, {len(df.columns)} columns):\n\n"
                    extracted_text += "Column Headers: " + ", ".join(df.columns) + "\n\n"
                    extracted_text += df.to_string(max_rows=100)  # Limit to first 100 rows
                    
                elif file_extension == '.xlsx':
                    # Handle modern Excel files
                    workbook = openpyxl.load_workbook(file_path, data_only=True)
                    text_lines = []
                    
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        text_lines.append(f"Sheet: {sheet_name}")
                        
                        # Extract text from cells
                        for row in sheet.iter_rows(values_only=True):
                            row_text = []
                            for cell_value in row:
                                if cell_value is not None:
                                    row_text.append(str(cell_value))
                            if row_text:
                                text_lines.append("\t".join(row_text))
                        text_lines.append("")
                    
                    extracted_text = '\n'.join(text_lines)
                    
                elif file_extension == '.xls':
                    # Handle older Excel files
                    workbook = xlrd.open_workbook(file_path)
                    text_lines = []
                    
                    for sheet_idx in range(workbook.nsheets):
                        sheet = workbook.sheet_by_index(sheet_idx)
                        text_lines.append(f"Sheet: {sheet.name}")
                        
                        for row_idx in range(sheet.nrows):
                            row_text = []
                            for col_idx in range(sheet.ncols):
                                cell_value = sheet.cell_value(row_idx, col_idx)
                                if cell_value:
                                    row_text.append(str(cell_value))
                            if row_text:
                                text_lines.append("\t".join(row_text))
                        text_lines.append("")
                    
                    extracted_text = '\n'.join(text_lines)
                    
            except Exception as e:
                logger.error(f"Excel/CSV extraction error: {str(e)}")
                extracted_text = f"Error processing {file_extension.upper()} file: {str(e)}"
            
            languages = [detect(extracted_text)] if extracted_text.strip() else ['en']
            
        elif file_extension in ['.txt', '.rtf']:
            # Process plain text files
            extracted_images = []
            
            try:
                with open(file_path, 'r', encoding='utf-8') as file:
                    extracted_text = file.read()
            except UnicodeDecodeError:
                # Try with different encoding
                try:
                    with open(file_path, 'r', encoding='latin-1') as file:
                        extracted_text = file.read()
                except Exception as e:
                    logger.error(f"Text file reading error: {str(e)}")
                    extracted_text = f"Error reading text file: {str(e)}"
            except Exception as e:
                logger.error(f"Text file processing error: {str(e)}")
                extracted_text = f"Error processing text file: {str(e)}"
            
            languages = [detect(extracted_text)] if extracted_text.strip() else ['en']
            
        else:
            # Fallback to original text extraction
            extracted_text = doc_processor.extract_text_from_file(file_path, mime_type)
            extracted_images = []
            languages = [detect(extracted_text)] if extracted_text.strip() else ['en']
        
        # If no text was extracted (common for images with poor OCR), use a safe fallback
        if not extracted_text or not extracted_text.strip():
            logger.warning("No text content found - using fallback text for analysis")
            if file_extension in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
                extracted_text = "Image document with no OCR text available. Analyze the visual content contextually for document type, sensitivity, and audience roles."
            else:
                extracted_text = "Document content not directly extractable. Provide general analysis based on filename and metadata."
            if not languages:
                languages = ['en']
        
        logger.info(f"Extracted {len(extracted_text)} characters and {len(extracted_images)} images")
        logger.info(f"Languages detected: {languages}")
        
        # Generate AI analysis
        analysis = processor.generate_ai_analysis(extracted_text)
        logger.info(f"AI analysis before normalization: {analysis.get('recommended_roles', {}).get('confidence', 'MISSING')}")

        # Normalize analysis to guarantee recommended_roles.confidence always exists and is sane
        analysis = normalize_analysis(analysis)
        logger.info(f"AI analysis after normalization: confidence={analysis.get('recommended_roles', {}).get('confidence', 'MISSING')}")

        # Add multilingual information to analysis
        analysis['languages_detected'] = languages
        analysis['images_extracted'] = len(extracted_images)
        analysis['has_multilingual_content'] = len(languages) > 1
        
        # Generate Malayalam summary if Malayalam content detected
        if 'ml' in languages and analysis.get('summary_en'):
            try:
                analysis['summary_ml'] = doc_processor.translate_text(analysis['summary_en'], 'ml')
            except Exception as e:
                logger.warning(f"Malayalam translation failed: {str(e)}")
                analysis['summary_ml'] = ""
        
        # Generate embeddings
        embeddings = processor.generate_embeddings(extracted_text)
        
        # Combine results
        result = {
            'extracted_text': extracted_text,
            'analysis': analysis,
            'embeddings': embeddings,
            'images': [{'id': img['id'], 'page': img['page'], 'text_content': img['text_content']} 
                      for img in extracted_images],  # Don't return full image data in main response
            'languages': languages,
            'processing_timestamp': datetime.now().isoformat()
        }
        
        return jsonify(result)
        
    except Exception as e:
        logger.error(f"Document processing error: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': str(e)}), 500

@app.route('/chat/document', methods=['POST'])
def chat_with_document():
    """Enhanced chat with entire document content using intelligent context selection."""
    try:
        data = request.json
        logger.info(f"Received chat request: {data}")
        
        document_text = data.get('document_text')
        user_question = data.get('question')
        language = data.get('language', 'en')
        document_id = data.get('document_id')
        
        # Enhanced validation with logging
        if not document_text:
            logger.error(f"Missing document_text in chat request. Data: {data}")
            return jsonify({'error': 'document_text is required and cannot be empty'}), 400
            
        if not user_question:
            logger.error(f"Missing question in chat request. Data: {data}")
            return jsonify({'error': 'question is required and cannot be empty'}), 400
            
        if len(document_text.strip()) == 0:
            logger.error(f"Empty document_text in chat request. Document ID: {document_id}")
            return jsonify({'error': 'Document text is empty. Please ensure the document has been processed successfully.'}), 400
        
        logger.info(f"Enhanced document chat - Document: {len(document_text)} chars, Question: {user_question[:100]}...")
        
        # Find relevant images for the question
        relevant_images = []
        if document_id:
            relevant_images = doc_processor.find_relevant_images(document_id, user_question, language)
        
        # Translate question if needed for better processing
        question_en = user_question if language == 'en' else doc_processor.translate_text(user_question, 'en')
        question_ml = user_question if language == 'ml' else doc_processor.translate_text(user_question, 'ml')
        
        # Create intelligent chunks from entire document
        chunks = chunk_processor.create_intelligent_chunks(document_text)
        logger.info(f"Created {len(chunks)} chunks for document chat")
        
        # Find most relevant chunks for the question
        relevant_chunks = []
        question_keywords = set(question_en.lower().split() + question_ml.lower().split())
        
        for chunk in chunks:
            # Calculate relevance score based on keyword matching
            chunk_words = set(chunk['text'].lower().split())
            keyword_matches = len(question_keywords.intersection(chunk_words))
            
            # Combine keyword relevance with section importance
            relevance_score = (keyword_matches / max(len(question_keywords), 1)) * 0.7 + chunk['importance'] * 0.3
            
            if relevance_score > 0.1 or any(keyword in chunk['text'].lower() for keyword in question_keywords):
                relevant_chunks.append({
                    'chunk': chunk,
                    'relevance_score': relevance_score
                })
        
        # Sort by relevance and take top chunks
        relevant_chunks.sort(key=lambda x: x['relevance_score'], reverse=True)
        top_chunks = relevant_chunks[:5]  # Use top 5 most relevant chunks
        
        logger.info(f"Found {len(top_chunks)} relevant chunks for question")
        
        # Create enhanced context from relevant chunks
        context_parts = []
        total_context_length = 0
        max_context_length = 8000  # Increased context limit
        
        for i, chunk_data in enumerate(top_chunks):
            chunk = chunk_data['chunk']
            chunk_text = chunk['text']
            
            # Add section headers if available
            section_header = ""
            if chunk['contains_sections']:
                section_header = f"[Section: {', '.join(chunk['contains_sections'])}]\n"
            
            chunk_context = f"{section_header}{chunk_text}"
            
            # Check if adding this chunk would exceed context limit
            if total_context_length + len(chunk_context) > max_context_length:
                # Truncate the chunk to fit
                remaining_space = max_context_length - total_context_length
                if remaining_space > 500:  # Only add if meaningful space remains
                    chunk_context = chunk_context[:remaining_space] + "..."
                    context_parts.append(f"--- Relevant Section {i+1} (Score: {chunk_data['relevance_score']:.2f}) ---\n{chunk_context}")
                break
            
            context_parts.append(f"--- Relevant Section {i+1} (Score: {chunk_data['relevance_score']:.2f}) ---\n{chunk_context}")
            total_context_length += len(chunk_context)
        
        # Combine document context
        document_context = "\n\n".join(context_parts)
        
        # Create image context
        image_context = ""
        if relevant_images:
            image_context = "\n\nRelevant Images Found:\n"
            for img in relevant_images[:3]:  # Limit to top 3 images
                image_context += f"- Image {img['id']} (Page {img['page']}): {img['text_content'][:200]}...\n"
        
        # Create enhanced prompt with full context
        prompt = f"""
        You are an AI assistant helping users understand documents for Kochi Metro Rail Limited (KMRL).
        You have access to the most relevant sections of the document for answering the user's question.
        
        Document Context (from {len(chunks)} total sections, showing {len(top_chunks)} most relevant):
        {document_context}
        {image_context}
        
        User Question (English): {question_en}
        User Question (Malayalam): {question_ml}
        
        Instructions:
        - Provide a comprehensive answer based on ALL the relevant document sections shown above
        - Reference specific sections when applicable
        - If information spans multiple sections, synthesize it coherently
        - Include relevant image information if available
        - Response language: {'Malayalam' if language == 'ml' else 'English'}
        - If the question cannot be fully answered from the provided sections, mention what information is available and what might be missing
        
        Provide a detailed, accurate answer based on the comprehensive document context.
        """
        
        response_text = optimized_deepseek_call(
            "enhanced_document_chat",
            f"{user_question[:100]}_{len(document_text)}",  # Cache key
            prompt,
            {'max_tokens': 3000, 'temperature': 0.3}
        )
        
        if not response_text:
            response_text = "I apologize, but I'm unable to process your question at the moment. Please try again later."
            if language == 'ml':
                response_text = "ക്ഷമിക്കണം, ഇപ്പോൾ നിങ്ങളുടെ ചോദ്യം പ്രോസസ്സ് ചെയ്യാൻ എനിക്ക് കഴിയുന്നില്ല. ദയവായി പിന്നീട് വീണ്ടും ശ്രമിക്കുക."
        
        # Prepare enhanced response with metadata
        response_data = {
            'answer': response_text,
            'language': language,
            'relevant_images': [{'id': img['id'], 'page': img['page'], 'relevance_score': img['relevance_score']} 
                               for img in relevant_images],
            'context_metadata': {
                'total_chunks': len(chunks),
                'relevant_chunks_used': len(top_chunks),
                'context_coverage_chars': total_context_length,
                'document_total_chars': len(document_text),
                'coverage_percentage': round((total_context_length / len(document_text)) * 100, 1),
                'processing_method': 'intelligent_context_selection'
            },
            'timestamp': datetime.now().isoformat()
        }
        
        return jsonify(response_data)
        
    except Exception as e:
        logger.error(f"Enhanced document chat error: {str(e)}")
        return jsonify({'error': str(e)}), 500
        
    except Exception as e:
        logger.error(f"Document chat error: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/images/<document_id>/<image_id>', methods=['GET'])
def get_document_image(document_id, image_id):
    """Get a specific image from a document"""
    try:
        # Ensure document_id is string for consistent lookup
        document_id = str(document_id)
        logger.info(f"Getting image {image_id} from document {document_id}")
        images = doc_processor.get_document_images(document_id)
        
        for image in images:
            if image['id'] == image_id:
                # Decode base64 image data
                image_data = base64.b64decode(image['data'])
                
                # Return image data with proper headers
                response = make_response(image_data)
                response.headers.set('Content-Type', f"image/{image.get('format', 'png')}")
                response.headers.set('Content-Length', len(image_data))
                response.headers.set('Cache-Control', 'max-age=3600')
                return response
        
        return jsonify({'error': 'Image not found'}), 404
        
    except Exception as e:
        logger.error(f"Image retrieval error: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/documents/<document_id>/images', methods=['GET'])
def get_document_images_list(document_id):
    """Get list of all images in a document"""
    try:
        # Ensure document_id is string for consistent lookup
        document_id = str(document_id)
        logger.info(f"Getting image list for document {document_id}")
        images = doc_processor.get_document_images(document_id)
        
        # Return image metadata without the actual image data
        image_list = [{
            'id': img['id'],
            'page': img['page'],
            'format': img['format'],
            'text_content': img['text_content'],
            'languages': img['languages'],
            'url': f"/images/{document_id}/{img['id']}"
        } for img in images]
        
        return jsonify({
            'document_id': document_id,
            'images': image_list,
            'total_images': len(image_list)
        })
        
    except Exception as e:
        logger.error(f"Image list error: {str(e)}")
        return jsonify({'error': str(e)}), 500

        return jsonify({'error': str(e)}), 500

# Removed duplicate global_chat function - keeping the more comprehensive one at line 1247

@app.route('/search/semantic', methods=['POST'])
def semantic_search():
    """Perform semantic search across documents."""
    try:
        data = request.json
        query = data.get('query')
        user_roles = data.get('user_roles', [])
        limit = data.get('limit', 10)
        
        if not query:
            return jsonify({'error': 'query is required'}), 400
        
        # Connect to database
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # Get documents accessible to user
        role_filter = "AND allowed_roles ?| %s" if user_roles else ""
        
        cursor.execute(f"""
            SELECT id, filename, summary_en, summary_ml, tags, embeddings, allowed_roles
            FROM documents 
            WHERE status = 'ACTIVE' {role_filter}
            ORDER BY created_at DESC
            LIMIT %s
        """, (user_roles, limit) if user_roles else (limit,))
        
        documents = cursor.fetchall()
        
        if not documents:
            return jsonify({'results': []})
        
        # Generate query embedding
        query_embedding = processor.generate_embeddings(query)
        
        # Calculate similarities
        results = []
        for doc in documents:
            doc_embeddings = doc['embeddings']
            if doc_embeddings:
                # Calculate cosine similarity
                similarity = cosine_similarity(
                    [query_embedding],
                    [doc_embeddings]
                )[0][0]
                
                results.append({
                    'document_id': doc['id'],
                    'filename': doc['filename'],
                    'summary_en': doc['summary_en'],
                    'summary_ml': doc['summary_ml'],
                    'tags': doc['tags'],
                    'similarity_score': float(similarity),
                    'allowed_roles': doc['allowed_roles']
                })
        
        # Sort by similarity
        results.sort(key=lambda x: x['similarity_score'], reverse=True)
        
        conn.close()
        
        return jsonify({
            'results': results[:limit],
            'query': query,
            'timestamp': datetime.now().isoformat()
        })
        
    except Exception as e:
        logger.error(f"Semantic search error: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/chat/global', methods=['POST'])
def global_chat():
    """Enhanced global chat with full document content access using intelligent chunking."""
    try:
        data = request.json
        user_question = data.get('query') or data.get('question')  # Support both field names
        user_roles = data.get('user_roles', [])
        language = data.get('language', 'en')
        
        if not user_question:
            return jsonify({'error': 'question is required'}), 400
        
        logger.info(f"Enhanced global chat query: {user_question[:100]}... Language: {language}")
        
        # Connect to database
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # Get documents accessible to user with more comprehensive filtering
        role_filter = ""
        params = []
        
        if user_roles:
            # Create role filter for JSON array
            role_conditions = []
            for role in user_roles:
                role_conditions.append("allowed_roles::text LIKE %s")
                params.append(f'%"{role}"%')
            role_filter = f"AND ({' OR '.join(role_conditions)})"
        
        # Get all relevant documents with full content for comprehensive analysis
        query = f"""
            SELECT id, filename, content, summary_en, summary_ml, tags, embeddings, 
                   allowed_roles, sensitivity_level, status,
                   created_at, file_size
            FROM documents 
            WHERE status = 'ACTIVE' {role_filter}
            ORDER BY created_at DESC
            LIMIT 20
        """
        
        cursor.execute(query, params)
        documents = cursor.fetchall()
        
        if not documents:
            return jsonify({
                'answer': 'No accessible documents found for your role.' if language == 'en' 
                         else 'നിങ്ങളുടെ റോളിന് ആക്സസ് ചെയ്യാവുന്ന ഡോക്യുമെന്റുകൾ കണ്ടെത്തിയില്ല.',
                'relevant_documents': [],
                'language': language,
                'timestamp': datetime.now().isoformat()
            })
        
        # Enhanced semantic search for relevance with full content analysis
        processor = DocumentProcessor()
        query_embedding = processor.generate_embeddings(user_question)
        
        relevant_docs = []
        document_chunks_analysis = []
        
        # Check if this is a general listing query
        listing_keywords = ['list', 'show', 'all documents', 'what documents', 'available documents', 'document list']
        is_listing_query = any(keyword in user_question.lower() for keyword in listing_keywords)
        
        # Translate question for better keyword matching
        question_en = user_question if language == 'en' else processor.translate_text(user_question, 'en')
        question_ml = user_question if language == 'ml' else processor.translate_text(user_question, 'ml')
        question_keywords = set((question_en + ' ' + question_ml).lower().split())
        
        for doc in documents:
            doc_relevance_data = {
                'id': doc['id'],
                'filename': doc['filename'],
                'summary_en': doc['summary_en'],
                'summary_ml': doc['summary_ml'],
                'tags': doc['tags'] or [],
                'allowed_roles': doc['allowed_roles'],
                'sensitivity_level': doc['sensitivity_level'],
                'status': doc['status'],
                'created_at': doc['created_at'].isoformat() if doc['created_at'] else None,
                'file_size': doc['file_size'],
                'similarity_score': 0.0,
                'content_matches': []
            }
            
            if doc['embeddings']:
                try:
                    doc_embedding = np.array(json.loads(doc['embeddings']))
                    similarity = cosine_similarity([query_embedding], [doc_embedding])[0][0]
                    doc_relevance_data['similarity_score'] = float(similarity)
                except:
                    similarity = 0.0
            else:
                similarity = 0.0
            
            # For non-listing queries, analyze full document content for better relevance
            if not is_listing_query and doc['content'] and similarity > 0.05:
                try:
                    # Create chunks from the full document content
                    chunks = chunk_processor.create_intelligent_chunks(doc['content'])
                    
                    # Find relevant chunks within this document
                    relevant_chunks = []
                    for chunk in chunks:
                        chunk_words = set(chunk['text'].lower().split())
                        keyword_matches = len(question_keywords.intersection(chunk_words))
                        
                        if keyword_matches > 0 or any(keyword in chunk['text'].lower() for keyword in question_keywords):
                            relevance_score = (keyword_matches / max(len(question_keywords), 1)) * 0.7 + chunk['importance'] * 0.3
                            relevant_chunks.append({
                                'text': chunk['text'][:500] + "..." if len(chunk['text']) > 500 else chunk['text'],
                                'relevance_score': relevance_score,
                                'section': chunk['contains_sections'][0] if chunk['contains_sections'] else 'Content'
                            })
                    
                    # Sort by relevance and take top 3 chunks per document
                    relevant_chunks.sort(key=lambda x: x['relevance_score'], reverse=True)
                    doc_relevance_data['content_matches'] = relevant_chunks[:3]
                    
                    # Boost document relevance if content matches found
                    if relevant_chunks:
                        max_chunk_relevance = max(chunk['relevance_score'] for chunk in relevant_chunks)
                        doc_relevance_data['similarity_score'] = max(similarity, max_chunk_relevance)
                    
                    document_chunks_analysis.append(doc_relevance_data)
                    
                except Exception as chunk_error:
                    logger.warning(f"Error chunking document {doc['filename']}: {str(chunk_error)}")
                    # Fallback to summary-based relevance
                    pass
            
            # Include document based on relevance thresholds
            threshold = 0.05 if is_listing_query else 0.15
            if doc_relevance_data['similarity_score'] > threshold or is_listing_query:
                relevant_docs.append(doc_relevance_data)
        
        # Sort by enhanced relevance score
        relevant_docs.sort(key=lambda x: x['similarity_score'], reverse=True)
        
        # Limit to top 10 most relevant documents for context
        top_docs = relevant_docs[:10]
        
        # Create comprehensive context with full document content
        context_parts = []
        document_stats = {
            'total_documents': len(documents),
            'relevant_documents': len(relevant_docs),
            'documents_with_content_analysis': len(document_chunks_analysis),
            'sensitivity_levels': {},
            'document_types': {},
            'recent_documents': 0
        }
        
        # Analyze document statistics
        from datetime import datetime, timedelta
        recent_cutoff = datetime.now() - timedelta(days=30)
        
        for doc in documents:
            # Count sensitivity levels
            sens_level = doc['sensitivity_level'] or 'UNKNOWN'
            document_stats['sensitivity_levels'][sens_level] = document_stats['sensitivity_levels'].get(sens_level, 0) + 1
            
            # Count document types (use generic types based on filename)
            filename = doc['filename'].lower()
            if any(ext in filename for ext in ['.pdf', '.doc', '.docx']):
                doc_type = 'document'
            elif any(ext in filename for ext in ['.xls', '.xlsx', '.csv']):
                doc_type = 'spreadsheet'
            elif any(ext in filename for ext in ['.ppt', '.pptx']):
                doc_type = 'presentation'
            else:
                doc_type = 'other'
            
            document_stats['document_types'][doc_type] = document_stats['document_types'].get(doc_type, 0) + 1
            
            # Count recent documents
            if doc['created_at'] and doc['created_at'] > recent_cutoff:
                document_stats['recent_documents'] += 1
        
        # Build enhanced context with content matches and collect relevant images
        all_relevant_images = []
        
        for i, doc in enumerate(top_docs[:5]):  # Include top 5 in context
            context_section = f"""
Document {i+1}: {doc['filename']} (Relevance: {doc['similarity_score']:.3f})
Status: {doc['status']}
Sensitivity: {doc['sensitivity_level']}
Summary: {doc['summary_en']}
Tags: {', '.join(doc['tags'][:5])}"""
            
            # Add relevant content matches if available
            if doc['content_matches']:
                context_section += "\nRelevant Content Sections:"
                for j, match in enumerate(doc['content_matches'][:2]):  # Top 2 matches per doc
                    context_section += f"\n  - [{match['section']}]: {match['text']}"
            
            # Find relevant images for each document
            try:
                doc_images = doc_processor.find_relevant_images(str(doc['id']), user_question, language)
                for img in doc_images[:2]:  # Top 2 images per document
                    img['document_id'] = doc['id']
                    img['document_filename'] = doc['filename']
                    all_relevant_images.append(img)
            except Exception as img_error:
                logger.warning(f"Error finding images for document {doc['filename']}: {str(img_error)}")
            
            context_parts.append(context_section.strip())
        
        # Sort all images by relevance and limit to top 5
        all_relevant_images.sort(key=lambda x: x.get('relevance_score', 0), reverse=True)
        top_images = all_relevant_images[:5]
        
        context = "\n\n".join(context_parts)
        
        # Add image context to the prompt if images are found
        image_context = ""
        if top_images:
            image_context = f"\n\nRelevant Images Found Across Documents:\n"
            for img in top_images:
                image_context += f"- Image {img['id']} from {img['document_filename']} (Page {img['page']}): {img['text_content'][:150]}...\n"
        
        # Enhanced prompt for comprehensive analysis with full content
        stats_summary = f"""
Total accessible documents: {document_stats['total_documents']}
Relevant to query: {document_stats['relevant_documents']}
Documents with full content analysis: {document_stats['documents_with_content_analysis']}
Recent documents (last 30 days): {document_stats['recent_documents']}
Sensitivity distribution: {', '.join([f"{k}: {v}" for k, v in document_stats['sensitivity_levels'].items()])}
Document types: {', '.join([f"{k}: {v}" for k, v in document_stats['document_types'].items()])}
        """
        
        prompt = f"""
You are an advanced AI assistant for Kochi Metro Rail Limited (KMRL) document management system.
You have access to {document_stats['total_documents']} documents with comprehensive content analysis for {document_stats['documents_with_content_analysis']} documents.

DOCUMENT COLLECTION OVERVIEW:
{stats_summary}

MOST RELEVANT DOCUMENTS WITH CONTENT ANALYSIS:
{context}
{image_context}

USER QUESTION: {user_question}

INSTRUCTIONS:
1. Provide a comprehensive answer based on the full document content and analysis
2. Reference specific documents and content sections when relevant
3. If analyzing trends, provide insights across multiple documents using the content matches
4. If the question is general, provide an overview using both summaries and actual content
5. Mention document types, sensitivity levels, or time periods when relevant
6. Use the relevant content sections to provide specific, accurate information
7. Include relevant image information if available and mention which documents they're from
8. Response language: {'Malayalam' if language == 'ml' else 'English'}
9. Be thorough but concise, highlighting key findings from the actual document content
10. If information spans multiple documents, synthesize it coherently
11. If information is not available in the analyzed content, clearly state this

Format your response in a clear, structured manner with specific references to document content and sections.
        """
        
        logger.info("Sending enhanced global query with full content analysis to DeepSeek V2 AI")
        
        response_text = optimized_deepseek_call(
            "enhanced_global_chat",
            f"{user_question[:50]}_{len(top_docs)}_{document_stats['documents_with_content_analysis']}",  # Cache key
            prompt,
            {'max_tokens': 4000, 'temperature': 0.3}
        )
        
        if not response_text:
            # Enhanced fallback response with content information
            fallback_answer = f"""
Based on comprehensive analysis of {document_stats['total_documents']} accessible documents with full content analysis on {document_stats['documents_with_content_analysis']} documents, I found {document_stats['relevant_documents']} documents relevant to your query.

Analysis Summary:
- Recent activity: {document_stats['recent_documents']} documents uploaded in the last 30 days
- Content matches found across {len([d for d in top_docs if d['content_matches']])} documents
- Document distribution: {', '.join([f"{k}: {v}" for k, v in document_stats['document_types'].items()])}

Most relevant documents for "{user_question}":
{chr(10).join([f"• {doc['filename']} (Relevance: {doc['similarity_score']:.3f}) - {len(doc['content_matches'])} content sections matched" for doc in top_docs[:5]])}

AI analysis temporarily unavailable. Please check individual documents for detailed information.
            """ if language == 'en' else f"""
{document_stats['documents_with_content_analysis']} ഡോക്യുമെന്റുകളുടെ സമ്പൂർണ്ണ ഉള്ളടക്ക വിശകലനത്തോടെ {document_stats['total_documents']} ആക്സസ് ചെയ്യാവുന്ന ഡോക്യുമെന്റുകളുടെ സമഗ്ര വിശകലനത്തെ അടിസ്ഥാനമാക്കി, നിങ്ങളുടെ ക്വറിയുമായി ബന്ധപ്പെട്ട് {document_stats['relevant_documents']} പ്രസക്തമായ ഡോക്യുമെന്റുകൾ കണ്ടെത്തി.

വിശദമായ വിവരങ്ങൾക്ക് വ്യക്തിഗത ഡോക്യുമെന്റുകൾ പരിശോധിക്കുക.
            """
            response_text = fallback_answer
        
        conn.close()
        
        return jsonify({
            'answer': response_text,
            'relevant_documents': top_docs,
            'relevant_images': [{'id': img['id'], 'page': img['page'], 'document_id': img['document_id'], 
                               'document_filename': img['document_filename'], 'relevance_score': img['relevance_score']} 
                               for img in top_images],
            'document_stats': document_stats,
            'language': language,
            'processing_metadata': {
                'full_content_analysis': True,
                'documents_analyzed': document_stats['documents_with_content_analysis'],
                'total_content_matches': sum(len(doc['content_matches']) for doc in top_docs),
                'total_relevant_images': len(top_images),
                'processing_method': 'intelligent_global_analysis'
            },
            'timestamp': datetime.now().isoformat()
        })
        
    except Exception as e:
        logger.error(f"Enhanced global chat error: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/debug/migrate-document-ids', methods=['POST'])
def migrate_document_ids():
    """Migrate existing UUID-based document images to use numeric IDs"""
    try:
        data = request.json
        mapping = data.get('mapping', {})  # Expected format: {"12": "03b2c480-0382-4867-b9e6-cd49d3699fda"}
        
        if not mapping:
            return jsonify({'error': 'mapping parameter required with format {"numeric_id": "uuid"}'}), 400
        
        results = {}
        
        for numeric_id, uuid_id in mapping.items():
            try:
                # Get images from UUID cache
                if uuid_id in doc_processor.extracted_images:
                    images = doc_processor.extracted_images[uuid_id]
                    
                    # Update image URLs to use numeric ID
                    updated_images = []
                    for image in images:
                        updated_image = image.copy()
                        # Update image ID to use numeric document ID
                        old_id = updated_image['id']
                        new_id = old_id.replace(uuid_id, numeric_id)
                        updated_image['id'] = new_id
                        updated_images.append(updated_image)
                    
                    # Store under numeric ID
                    doc_processor.extracted_images[numeric_id] = updated_images
                    
                    # Save to database if possible
                    doc_processor.save_images_to_database(numeric_id, updated_images)
                    
                    results[numeric_id] = {
                        'status': 'success',
                        'images_migrated': len(updated_images),
                        'source_uuid': uuid_id
                    }
                    
                    logger.info(f"Migrated {len(updated_images)} images from UUID {uuid_id} to numeric ID {numeric_id}")
                else:
                    results[numeric_id] = {
                        'status': 'not_found',
                        'message': f'No images found for UUID {uuid_id}'
                    }
                    
            except Exception as e:
                results[numeric_id] = {
                    'status': 'error',
                    'message': str(e)
                }
        
        return jsonify({
            'message': 'Migration completed',
            'results': results
        })
        
    except Exception as e:
        logger.error(f"Migration error: {str(e)}")
        return jsonify({'error': f'Migration failed: {str(e)}'}), 500

@app.route('/debug/images-cache', methods=['GET'])
def debug_images_cache():
    """Debug endpoint to check what's in the images cache"""
    try:
        cache_info = {}
        for doc_id, images in doc_processor.extracted_images.items():
            cache_info[doc_id] = {
                'image_count': len(images),
                'image_ids': [img.get('id', 'unknown') for img in images[:3]]  # First 3 image IDs
            }
        
        return jsonify({
            'total_documents_in_cache': len(doc_processor.extracted_images),
            'cache_details': cache_info,
            'cache_keys': list(doc_processor.extracted_images.keys())
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/translate', methods=['POST'])
def translate_text():
    """Translate text between English and Malayalam"""
    try:
        data = request.json
        text = data.get('text')
        target_language = data.get('target_language', 'en')
        
        if not text:
            return jsonify({'error': 'text is required'}), 400
        
        if target_language not in ['en', 'ml']:
            return jsonify({'error': 'target_language must be en or ml'}), 400
        
        # Use the document processor's translation capability
        translated_text = doc_processor.translate_text(text, target_language)
        
        if not translated_text:
            translated_text = text  # Fallback to original text
        
        return jsonify({
            'translated_text': translated_text,
            'source_language': 'en' if target_language == 'ml' else 'ml',
            'target_language': target_language,
            'timestamp': datetime.now().isoformat()
        })
        
    except Exception as e:
        logger.error(f"Translation error: {str(e)}")
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(
        host='0.0.0.0',
        port=int(os.getenv('FLASK_PORT', 5000)),
        debug=True
    )
